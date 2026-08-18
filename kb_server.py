# -*- coding: utf-8 -*-
"""
MINISFORUM KB Server V10 — FastAPI + SSE Streaming + Cache
- Migrated from V9's BaseHTTPRequestHandler to FastAPI + uvicorn
- SSE streaming for /api/kb/ask and /api/chat
- In-memory query cache with TTL
- Async-capable architecture
"""

import os, re, json, sys, io, hashlib, time, sqlite3, csv, html, subprocess, tempfile, shutil, base64, math, threading, logging, heapq, asyncio
import numpy as np
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime

# ── Logger ──────────────────────────────────────────────
logging.basicConfig(level=logging.INFO, format='[%(asctime)s] %(name)s %(levelname)s: %(message)s')
logger = logging.getLogger("kb_v10")

# ── Config ──────────────────────────────────────────────
NAS_DOCS_DIR  = os.environ.get("KB_DOCS_DIR", "/kb_persist/docs")
DB_PATH       = os.environ.get("KB_DB_PATH", "/kb_persist/kb.db")
OLLAMA_HOST   = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
# Compare module uses a separate GPU-accelerated Ollama (14b on Radeon 8060S)
OLLAMA_HOST_COMPARE = os.environ.get("KB_OLLAMA_HOST_COMPARE", "http://localhost:11434")

# ── MinisLLM (35B via faster_llm OpenAI-compatible API, non-ollama) ──
MINISLLM_URL   = os.environ.get("KB_MINISLLM_URL", "http://localhost:30000/v1")
MINISLLM_MODEL = os.environ.get("KB_MINISLLM_MODEL", "qwen2.5:35b")

# ── Model Configuration (per-module) ──────────────
# Each feature module can use a different Ollama model.
# Set via env vars: module-specific → KB_MODEL (global fallback) → hardcoded default.
# doc compare / spec compare default to 14b because 7b produces poor comparison tables.
_DEFAULT = "qwen2.5:14b"
_GLOBAL  = os.environ.get("KB_MODEL", "")
KB_MODEL_ASK       = os.environ.get("KB_MODEL_ASK",       _GLOBAL or "qwen2.5:7b")
KB_MODEL_CHAT      = os.environ.get("KB_MODEL_CHAT",      _GLOBAL or _DEFAULT)
KB_MODEL_ASK_EN    = os.environ.get("KB_MODEL_ASK_EN",    _GLOBAL or "qwen2.5:7b")
KB_MODEL_COMPARE   = os.environ.get("KB_MODEL_COMPARE",   _GLOBAL or "qwen2.5:14b")
KB_MODEL_EXTRACTOR = os.environ.get("KB_MODEL_EXTRACTOR",
                      os.environ.get("KB_EXTRACTOR_MODEL", _GLOBAL or _DEFAULT))
# Backwards compat aliases
LLM_MODEL        = KB_MODEL_ASK
EXTRACTOR_MODEL  = KB_MODEL_EXTRACTOR

LLM_API_URL   = ""  # hardcoded
LLM_API_TOKEN = os.environ.get("KB_LLM_API_TOKEN", "")

# ── GPU 35B channel (faster_llm on Radeon 8060S, OpenAI-compatible) ──
GPU_LLM_URL   = os.environ.get("KB_GPU_LLM_URL", "http://localhost:13306/v1")
GPU_LLM_MODEL = os.environ.get("KB_GPU_LLM_MODEL", "Qwen_Qwen3.6-35B-A3B-Q4_K_M.gguf")
EMBED_MODEL   = os.environ.get("KB_EMBED_MODEL", "nomic-embed-text")
VECTOR_DIM    = int(os.environ.get("KB_VECTOR_DIM", "768"))  # nomic-embed-text
TOP_K         = int(os.environ.get("KB_TOP_K", "6"))
CHUNK_MAX     = int(os.environ.get("KB_CHUNK_MAX", "800"))
CHUNK_OVERLAP = int(os.environ.get("KB_CHUNK_OVERLAP", "100"))
MAX_CONTEXT_CHARS = int(os.environ.get("KB_MAX_CONTEXT_CHARS", "3000"))
JWT_SECRET    = os.environ.get("KB_JWT_SECRET", "kb_local_secret_change_me")
JWT_ALG       = "HS256"
CACHE_TTL     = int(os.environ.get("KB_CACHE_TTL", "300"))

# os.environ["OLLAMA_HOST"] = OLLAMA_HOST  # disabled: hardcoded above

# ── OCR (PaddleOCR, CPU) ────────────────────────────────
# Enable OCR for scanned PDFs / images on upload. Set KB_OCR_ENABLED=0 to disable.
KB_OCR_ENABLED = os.environ.get("KB_OCR_ENABLED", "1") == "1"
OCR_DPI       = int(os.environ.get("KB_OCR_DPI", "200"))

# ── Imports ─────────────────────────────────────────────
import urllib.request as _ur

# ── Query Cache ─────────────────────────────────────────
_cache = {}
_cache_lock = threading.Lock()

def cache_get(key):
    with _cache_lock:
        entry = _cache.get(key)
        if entry and time.time() - entry["ts"] < CACHE_TTL:
            return entry["value"]
        if entry:
            del _cache[key]
    return None

def cache_set(key, value):
    with _cache_lock:
        _cache[key] = {"value": value, "ts": time.time()}

def cache_invalidate():
    with _cache_lock:
        _cache.clear()
        _embed_cache.clear()
    global _emb_index
    _emb_index = None  # force rebuild on next vector_search

# ── Question Quality ────────────────────────────────────
def classify_question(question):
    """Returns True if question is too vague (keyword-only, no question intent)."""
    q = question.strip()
    question_words = [
        '什么', '怎么', '如何', '为什么', '谁', '哪', '吗', '呢', '多少',
        '介绍', '说明', '对比', '区别', '功能', '参数', '价格', '配置',
        '规格', '评测', '使用', '安装', '下载', '教程', '指南', '推荐',
        'what', 'how', 'why', 'which', 'where', 'when', 'who', 'tell',
        'explain', 'describe', 'show', 'list', 'compare', 'spec', 'price',
        'review', 'guide', 'tutorial', 'install', 'download', 'feature',
    ]
    has_indicator = any(w in q.lower() for w in question_words)
    has_mark = '?' in q or '？' in q
    if has_indicator or has_mark:
        return False
    if len(q) <= 15:
        return True
    chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', q))
    if chinese_chars < 3:
        return True
    return False

# ── Embedding Index (numpy-accelerated batch cosine similarity) ──
_emb_index = None       # {"ids": [(doc_id,chunk_index), ...], "texts": [text, ...], "matrix": np.ndarray, "norms": np.ndarray}
_emb_index_lock = threading.Lock()

def _build_embedding_index():
    """Pre-load all chunk embeddings into numpy matrix. Called lazily on first vector_search or after imports."""
    global _emb_index
    chunks = db_query("SELECT doc_id, chunk_index, text, embedding FROM chunks")
    if not chunks:
        _emb_index = None
        return
    
    ids, texts, vectors = [], [], []
    for row in chunks:
        emb_str = row["embedding"]
        if not emb_str:
            continue
        try:
            emb = json.loads(emb_str)
        except (json.JSONDecodeError, TypeError):
            continue
        if not emb or len(emb) == 0:
            continue
        ids.append((row["doc_id"], row["chunk_index"]))
        texts.append(row["text"] or "")
        vectors.append(emb)
    
    if not vectors:
        _emb_index = None
        return
    
    matrix = np.array(vectors, dtype=np.float32)
    norms = np.linalg.norm(matrix, axis=1)
    
    _emb_index = {"ids": ids, "texts": texts, "matrix": matrix, "norms": norms}
    print(f"[V10] Embedding index built: {len(ids)} chunks, shape {matrix.shape}", flush=True)

# ── Stats ───────────────────────────────────────────────
_stats = {"queries_today": 0, "today_date": datetime.now().strftime("%Y-%m-%d")}
_stats_lock = threading.Lock()

# ── Query Concurrency Control ───────────────────────────
# Only one LLM query (kb_ask / chat / compare) runs at a time; others queue up.
_query_semaphore = threading.Semaphore(1)
_query_queue_counter = 0
_query_counter_lock = threading.Lock()

def _incr_query():
    today = datetime.now().strftime("%Y-%m-%d")
    with _stats_lock:
        if _stats["today_date"] != today:
            _stats["today_date"] = today
            _stats["queries_today"] = 0
        _stats["queries_today"] += 1

# ── Ollama helpers ──────────────────────────────────────
_embed_cache: dict = {}

def ollama_embed(text):
    """调用 ollama API 做向量化（nomic-embed-text, 768维）"""
    key = hashlib.md5(text.encode()).hexdigest()
    if key in _embed_cache:
        return _embed_cache[key]
    url = f"{OLLAMA_HOST}/api/embeddings"
    data = json.dumps({"model": EMBED_MODEL, "prompt": text}).encode()
    req = _ur.Request(url, data=data, headers={"Content-Type": "application/json"})
    try:
        resp = _ur.urlopen(req, timeout=30)
        result = json.loads(resp.read())
        emb = result["embedding"]
        _embed_cache[key] = emb
        return emb
    except Exception as e:
        logger.error(f"[Embed] ollama API failed: {e}")
        raise

def ollama_chat(messages, stream=False, model=None, max_tokens=4096, temperature=0.3, base_url=None, use_gpu=False):
    if use_gpu:
        # GPU 35B channel via faster_llm (OpenAI-compatible). enable_thinking=false
        # to avoid reasoning model consuming all max_tokens and returning empty content.
        url = f"{GPU_LLM_URL}/chat/completions"
        data = {"model": GPU_LLM_MODEL, "messages": messages, "max_tokens": max_tokens,
                "stream": stream, "temperature": temperature,
                "chat_template_kwargs": {"enable_thinking": False}}
        headers = {"Content-Type": "application/json"}
        if LLM_API_TOKEN:
            headers["Authorization"] = f"Bearer {LLM_API_TOKEN}"
        req = _ur.Request(url, data=json.dumps(data).encode(), headers=headers)
        if stream:
            return _ur.urlopen(req, timeout=900)
        with _ur.urlopen(req, timeout=900) as resp:
            r = json.loads(resp.read())
            return r["choices"][0]["message"]["content"]
    if LLM_API_URL:
        url = f"{LLM_API_URL}/chat/completions"
        data = {"model": "default", "messages": messages, "max_tokens": max_tokens, "stream": stream, "temperature": temperature}
        headers = {"Content-Type": "application/json"}
        if LLM_API_TOKEN:
            headers["Authorization"] = f"Bearer {LLM_API_TOKEN}"
        req = _ur.Request(url, data=json.dumps(data).encode(), headers=headers)
        if stream:
            return _ur.urlopen(req, timeout=900)
        with _ur.urlopen(req, timeout=900) as resp:
            r = json.loads(resp.read())
            return r["choices"][0]["message"]["content"]
    else:
        model = model or LLM_MODEL
        data = {"model": model, "messages": messages, "stream": stream, "keep_alive": "30m", "options": {"temperature": temperature, "num_predict": max_tokens, "num_ctx": 16384}}
        url = f"{base_url or OLLAMA_HOST}/api/chat"
        req = _ur.Request(url, data=json.dumps(data).encode(), headers={"Content-Type": "application/json"})
        if stream:
            return _ur.urlopen(req, timeout=900)
        with _ur.urlopen(req, timeout=900) as resp:
            r = json.loads(resp.read())
            return r["message"]["content"]

def ollama_chat_stream(messages, model=None, max_tokens=4096, temperature=0.3, use_gpu=False):
    """Generator that yields content deltas from streaming LLM response.
    Handles three formats:
    - SSE (faster_llm OpenAI compat): `data: {"choices":[{"delta":{"content":"..."}}]}\n\n`
    - SSE (Native Ollama wrapped): `data: {"message":{"content":"..."}}\n\n`
    - NDJSON (Native Ollama raw): `{"model":"...","message":{"content":"..."}}\n`
    """
    import logging, time
    _log = logging.getLogger("kb_stream")
    t0 = time.time()
    resp = ollama_chat(messages, stream=True, model=model, max_tokens=max_tokens, temperature=temperature, use_gpu=use_gpu)
    n_lines = n_reasoning = n_content = 0
    first_content_ts = None
    for line_bytes in resp:
        n_lines += 1
        line = line_bytes.rstrip(b"\r\n")
        if not line:
            continue
        if line == b"data: [DONE]":
            _log.info("[stream] done after %d lines | reasoning=%d content=%d | %.1fs",
                      n_lines, n_reasoning, n_content, time.time() - t0)
            return
        # Try SSE "data: {...}" first, then raw NDJSON "{...}" for Native Ollama
        json_str = None
        if line.startswith(b"data: "):
            json_str = line[6:]
        elif line.startswith(b"{"):
            json_str = line
        else:
            continue
        try:
            obj = json.loads(json_str)
            if "choices" in obj:
                delta = obj.get("choices", [{}])[0].get("delta", {})
                reasoning = delta.get("reasoning_content", "")
                content = delta.get("content", "")
            elif "message" in obj:
                reasoning = ""
                content = obj.get("message", {}).get("content", "")
            else:
                continue
            # Guard against JSON null (None) masquerading as content
            if content is None:
                content = ""
            if reasoning:
                n_reasoning += 1
                yield ("reasoning", reasoning)
            if content:
                if first_content_ts is None:
                    first_content_ts = time.time() - t0
                    _log.info("[stream] first CONTENT at %.2fs (line %d)", first_content_ts, n_lines)
                n_content += 1
                yield ("token", content)
        except Exception:
            pass

def minisllm_chat(messages, max_tokens=4096, temperature=0.3):
    """Call MinisLLM 35B via OpenAI-compatible API (faster_llm wrapper, non-ollama)."""
    url = f"{MINISLLM_URL}/chat/completions"
    data = {"model": MINISLLM_MODEL, "messages": messages,
            "max_tokens": max_tokens, "temperature": temperature}
    req = _ur.Request(url, data=json.dumps(data).encode(),
                      headers={"Content-Type": "application/json"})
    with _ur.urlopen(req, timeout=900) as resp:
        r = json.loads(resp.read())
        return r["choices"][0]["message"]["content"]

def extract_facts(context_text, question):
    prompt = (
        "You are a fact extraction engine. The user wants to know: \"{question}\"\n"
        "Extract ALL facts from the documents below that help answer this specific question.\n"
        "Focus on what the user is actually asking — company history, products, data, events, etc.\n\n"
        "RULES:\n"
        "1. Extract explicit facts, numbers, specs, data points, dates, events, milestones, "
        "awards, achievements, and company/product descriptions.\n"
        "2. MATCHING RULE: Treat all case, hyphen, underscore, and space variations as the SAME product name.\n"
        "   - Normalize by removing all hyphens, underscores, and spaces, then lowercase, before comparing.\n"
        "   - Example: 'N5MAX', 'N5_MAX', 'N5-MAX', 'n5 max' all refer to the same product.\n"
        "3. Keep original values exactly as written (units, formatting, precision).\n"
        "4. Output one fact per line, as: '[Document: doc_name] | field: value'\n"
        "5. If truly nothing is relevant, output only: NO_RELEVANT_FACTS\n\n"
        f"User question: {question}\n\n"
        f"Documents:\n{context_text}"
    )
    return ollama_chat([
        {"role": "system", "content": "You are a fact extraction engine. Extract only explicit facts."},
        {"role": "user", "content": prompt}
    ], model=EXTRACTOR_MODEL, max_tokens=1024)

# ── DB ──────────────────────────────────────────────────
def db_exec(sql, params=None):
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    if params: cur.execute(sql, params)
    else: cur.execute(sql)
    conn.commit()
    conn.close()

def db_query(sql, params=None):
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    cur = conn.cursor()
    if params: cur.execute(sql, params)
    else: cur.execute(sql)
    rows = [dict(r) for r in cur.fetchall()]
    conn.close()
    return rows

def init_db():
    db_exec("""CREATE TABLE IF NOT EXISTS documents (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        title TEXT, filename TEXT, filepath TEXT, filetype TEXT,
        department TEXT DEFAULT 'general', upload_user TEXT,
        created_at DATETIME DEFAULT CURRENT_TIMESTAMP
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS chunks (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        doc_id INTEGER, chunk_index INTEGER, text TEXT,
        embedding BLOB, keywords TEXT, lang TEXT DEFAULT 'zh',
        FOREIGN KEY(doc_id) REFERENCES documents(id) ON DELETE CASCADE
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE, password_hash TEXT,
        department TEXT DEFAULT 'general', role TEXT DEFAULT 'user',
        status TEXT DEFAULT 'pending',
        created_at DATETIME DEFAULT CURRENT_TIMESTAMP
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS sessions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        session_id TEXT, role TEXT, content TEXT,
        created_at DATETIME DEFAULT CURRENT_TIMESTAMP
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS queries (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        query TEXT, session_id TEXT, answered INTEGER DEFAULT 0,
        created_at DATETIME DEFAULT CURRENT_TIMESTAMP
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS feedback (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        query_id INTEGER, rating TEXT,
        created_at DATETIME DEFAULT CURRENT_TIMESTAMP
    )""")
    db_exec("""CREATE TABLE IF NOT EXISTS doc_versions (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        group_id TEXT NOT NULL,
        doc_id INTEGER NOT NULL,
        version INTEGER NOT NULL DEFAULT 1,
        created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
        FOREIGN KEY(doc_id) REFERENCES documents(id) ON DELETE CASCADE
    )""")
    # Schema migration
    for tbl, col, typ in [("documents", "filepath", "TEXT"), ("documents", "filetype", "TEXT"), ("documents", "upload_user", "TEXT"), ("documents", "version_group", "TEXT"), ("documents", "version", "INTEGER DEFAULT 1")]:
        try:
            db_exec(f"ALTER TABLE {tbl} ADD COLUMN {col} {typ}")
        except sqlite3.OperationalError:
            pass
    try:
        db_exec("ALTER TABLE chunks ADD COLUMN lang TEXT DEFAULT 'zh'")
    except sqlite3.OperationalError:
        pass
    # Populate version_group for existing docs
    import uuid as _uuid
    orphan_docs = db_query("SELECT id FROM documents WHERE version_group IS NULL OR version IS NULL")
    for d in orphan_docs:
        db_exec("UPDATE documents SET version_group=?, version=1 WHERE id=?", (str(_uuid.uuid4()), d["id"]))
        db_exec("INSERT INTO doc_versions (group_id, doc_id, version) VALUES (?,?,1)",
                (db_query("SELECT version_group FROM documents WHERE id=?", (d["id"],))[0]["version_group"], d["id"]))
    # Ensure default admin
    existing = db_query("SELECT id FROM users WHERE username='admin'")
    if not existing:
        try:
            from passlib.hash import bcrypt as _bcrypt
            pwd_hash = _bcrypt.hash("admin123")
        except ImportError:
            raise RuntimeError("passlib is required for password hashing. Install with: pip install passlib[bcrypt]")
        db_exec("INSERT INTO users (username,password_hash,department,role,status) VALUES (?,?,?,?,?)",
                ("admin", pwd_hash, "Technology", "superadmin", "approved"))
        print("[V10] Default admin created", flush=True)
    # Startup security checks
    if not LLM_API_TOKEN:
        print("[V10] WARNING: LLM_API_TOKEN is using default value. Set KB_LLM_API_TOKEN env var.", flush=True)
    if JWT_SECRET == "kb_local_secret_change_me":
        print("[V10] WARNING: JWT_SECRET is using default value. Set KB_JWT_SECRET env var.", flush=True)

# ── BM25 ────────────────────────────────────────────────
class BM25:
    def __init__(self):
        self.k1, self.b = 1.5, 0.75
        self.doc_freqs, self.term_freqs, self.doc_lengths = {}, {}, {}
        self.total_docs, self.avgdl = 0, 0

    def tokenize(self, text):
        return re.findall(r'[\u4e00-\u9fff]+|[a-zA-Z0-9]+', text.lower())

    def add(self, doc_id, chunk_id, text):
        tokens = self.tokenize(text)
        key = (doc_id, chunk_id)
        self.doc_lengths[key] = len(tokens)
        self.total_docs += 1
        tf = {}
        for t in tokens:
            tf[t] = tf.get(t, 0) + 1
            self.doc_freqs[t] = self.doc_freqs.get(t, 0) + 1
        self.term_freqs[key] = tf
        self.avgdl = sum(self.doc_lengths.values()) / max(1, self.total_docs)

    def remove(self, doc_id):
        keys = [k for k in self.doc_lengths if k[0] == doc_id]
        for key in keys:
            if key in self.term_freqs:
                for t in self.term_freqs[key]:
                    self.doc_freqs[t] = max(0, self.doc_freqs.get(t, 1) - 1)
                del self.term_freqs[key]
            if key in self.doc_lengths:
                self.total_docs -= 1
                del self.doc_lengths[key]
        self.avgdl = sum(self.doc_lengths.values()) / max(1, self.total_docs)

    def search(self, query, top_k=10):
        tokens = self.tokenize(query)
        if not tokens or self.total_docs == 0:
            return []
        scores = {}
        for t in tokens:
            df = self.doc_freqs.get(t, 0)
            if df == 0:
                continue
            try:
                idf = max(0, math.log((self.total_docs - df + 0.5) / (df + 0.5)) + 1)
            except ValueError:
                idf = 0
            for key, tf_map in self.term_freqs.items():
                if t in tf_map:
                    f = tf_map[t]
                    dl = self.doc_lengths.get(key, 1)
                    score = idf * (f * (self.k1 + 1)) / (f + self.k1 * (1 - self.b + self.b * dl / max(1, self.avgdl)))
                    scores[key] = scores.get(key, 0) + score
        ranked = sorted(scores.items(), key=lambda x: -x[1])[:top_k]
        results = []
        for key, score in ranked:
            doc_id, chunk_id = key
            rows = db_query("SELECT text FROM chunks WHERE doc_id=? AND chunk_index=?", (doc_id, chunk_id))
            if rows:
                results.append((doc_id, chunk_id, rows[0]["text"], score))
        return results

    def rebuild(self):
        self.doc_freqs, self.term_freqs, self.doc_lengths = {}, {}, {}
        self.total_docs, self.avgdl = 0, 0
        for row in db_query("SELECT doc_id, chunk_index, text FROM chunks"):
            self.add(row["doc_id"], row["chunk_index"], row["text"])
        self.avgdl = sum(self.doc_lengths.values()) / max(1, self.total_docs)

bm25 = BM25()

# ── Vector helpers ──────────────────────────────────────
def vector_search(query, top_k=10, dept=None, role=None, lang=None):
    t0 = time.time()
    query_vec = ollama_embed(query)
    t_embed = time.time()

    # Lazy-build or refresh embedding index
    with _emb_index_lock:
        if _emb_index is None:
            _build_embedding_index()
        idx = _emb_index
        if idx is None:
            return []
        ids = idx["ids"]
        texts = idx["texts"]
        matrix = idx["matrix"]
        norms = idx["norms"]
    t_index = time.time()

    # Department filter
    need_dept_filter = bool(dept and role and role not in ("admin", "superadmin"))
    if need_dept_filter:
        dept_rows = db_query("SELECT id, department FROM documents")
        dept_map = {r["id"]: (r.get("department") or "") for r in dept_rows}
        mask = np.array([dept_map.get(doc_id, "") == dept for doc_id, _ in ids], dtype=bool)
        if not mask.any():
            return []
        matrix = matrix[mask]
        norms = norms[mask]
        ids = [ids[i] for i in range(len(ids)) if mask[i]]
        texts = [texts[i] for i in range(len(texts)) if mask[i]]
    t_filter = time.time()

    # Batch cosine similarity via numpy
    q_vec = np.array(query_vec, dtype=np.float32)
    q_norm = float(np.linalg.norm(q_vec))
    if q_norm == 0:
        return []
    dots = np.dot(matrix, q_vec)
    sims = dots / (q_norm * norms)
    t_sim = time.time()

    # Top-k
    k = min(top_k, len(sims))
    if k == 0:
        return []
    top_indices = np.argpartition(-sims, k - 1)[:k]
    top_indices = top_indices[np.argsort(-sims[top_indices])]
    t_topk = time.time()

    results = [((ids[i][0], ids[i][1]), {"text": texts[i], "vec_score": float(sims[i])}) for i in top_indices]

    total_ms = (t_topk - t0) * 1000
    if total_ms > 500:
        parts = f"embed={(t_embed-t0)*1000:.0f}ms index={(t_index-t_embed)*1000:.0f}ms filter={(t_filter-t_index)*1000:.0f}ms sim={(t_sim-t_filter)*1000:.0f}ms topk={(t_topk-t_sim)*1000:.0f}ms"
        logger.info("[Perf] vector_search total=%.0fms (%s) → %d results", total_ms, parts, len(results))
    return results

def _bm25_search_wrapper(query, top_k):
    return bm25.search(query, top_k)

def hybrid_search(query, top_k=10, dept=None, role=None, bm25_weight=1.0, lang=None):
    with ThreadPoolExecutor(max_workers=2) as executor:
        vec_future = executor.submit(vector_search, query, top_k * 3, dept, role, lang)
        bm25_future = executor.submit(_bm25_search_wrapper, query, top_k * 3)
        vec_results = vec_future.result()
        bm25_results = bm25_future.result()
    vec_map = {key: info for key, info in vec_results}
    bm25_map = {(doc_id, chunk_id): score for doc_id, chunk_id, _, score in bm25_results}
    all_keys = set(vec_map.keys()) | set(bm25_map.keys())
    k = 60
    combined = []
    for key in all_keys:
        v_rank = next((i + 1 for i, (k2, _) in enumerate(vec_results) if k2 == key), len(vec_results) + 1)
        b_rank = next((i + 1 for i, (d, c, _, _) in enumerate(bm25_results) if (d, c) == key), len(bm25_results) + 1)
        rrf = 1.0 / (k + v_rank) + bm25_weight / (k + b_rank)
        info = vec_map.get(key, {"text": "", "vec_score": None})
        info["rrf_score"] = rrf
        info["bm25_rank"] = b_rank
        combined.append((key, info))
    combined.sort(key=lambda x: -x[1]["rrf_score"])
    if lang is not None:
        # Language boost: same-lang chunks get +0.15 RRF score; different lang still included as fallback
        for k, v in combined:
            rows = db_query("SELECT lang FROM chunks WHERE doc_id=? AND chunk_index=?", (k[0], k[1]))
            cl = rows[0]["lang"] if rows else "zh"
            if cl == lang:
                v["rrf_score"] += 0.15
        combined.sort(key=lambda x: -x[1]["rrf_score"])
    logger.info("[Diag] hybrid_search top-%d (bm25=%d vec=%d combined=%d):", top_k, len(bm25_results), len(vec_results), len(combined))
    for i, (key, info) in enumerate(combined[:top_k]):
        title = get_document_title(key[0])
        chk_text = get_chunk_text(key[0], key[1])
        logger.info("[Diag]   #%d doc=%s chk=%d rrf=%.4f bm25_r=%d title=%s",
                    i+1, key[0], key[1], info["rrf_score"], info.get("bm25_rank", -1), title[:80])
        logger.info("[Diag]   chunk_%d_text: %s", i+1, chk_text[:120])
    return combined[:top_k]

def get_document_title(doc_id):
    rows = db_query("SELECT title FROM documents WHERE id=?", (doc_id,))
    return rows[0]["title"] if rows else str(doc_id)

def get_chunk_text(doc_id, chunk_idx):
    rows = db_query("SELECT text FROM chunks WHERE doc_id=? AND chunk_index=?", (doc_id, chunk_idx))
    return rows[0]["text"] if rows else ""

# ── Smart Chunking ──────────────────────────────────────
def smart_chunk(text):
    paragraphs = re.split(r'\n\s*\n', text)
    chunks, current = [], ""
    if not text.strip():
        logger.warning("[smart_chunk] Input text is empty/whitespace-only, returning 0 chunks")
        return chunks
    for para in paragraphs:
        para = para.strip()
        if not para: continue
        if len(current) + len(para) < CHUNK_MAX:
            current = (current + "\n\n" + para).strip()
        else:
            if current: chunks.append(current)
            current = para
        if len(current) >= CHUNK_MAX:
            sentences = re.split(r'(?<=[.!?。！？])\s+', current)
            current = ""
            for sent in sentences:
                if len(current) + len(sent) < CHUNK_MAX:
                    current = (current + " " + sent).strip()
                else:
                    if current: chunks.append(current)
                    current = sent
            if len(current) > CHUNK_MAX:
                for i in range(0, len(current), CHUNK_MAX):
                    chunks.append(current[i:i + CHUNK_MAX])
                current = ""
    if current: chunks.append(current)
    # Apply overlap: prepend tail of previous chunk to each subsequent chunk
    if CHUNK_OVERLAP > 0 and len(chunks) > 1:
        overlap_chunks = [chunks[0]]
        for i in range(1, len(chunks)):
            prev_tail = chunks[i-1][-CHUNK_OVERLAP:] if len(chunks[i-1]) >= CHUNK_OVERLAP else chunks[i-1]
            overlap_chunks.append(prev_tail + "\n" + chunks[i])
        return overlap_chunks
    return chunks

# ── Parsers ─────────────────────────────────────────────
def _parse_raw(filepath):
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            with open(filepath, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError: continue
    return ""

def _parse_txt(fp): return _parse_raw(fp)
def _parse_md(fp):
    text = _parse_raw(fp)
    # Strip YAML front matter to avoid metadata polluting chunks
    if text.startswith('---'):
        parts = text.split('---', 2)
        if len(parts) >= 3:
            text = parts[2].strip()
    return text
def _parse_html_file(fp):
    raw = _parse_raw(fp)
    from html.parser import HTMLParser
    class P(HTMLParser):
        def __init__(self):
            super().__init__()
            self.text = []
        def handle_data(self, d):
            d = d.strip()
            if d: self.text.append(d)
    p = P(); p.feed(raw)
    return "\n".join(p.text)

# ── OCR (PaddleOCR, CPU) ─────────────────────────────
_ocr_engine = None
_ocr_lock = threading.Lock()

def _get_ocr_engine():
    """Lazy-load PaddleOCR engine (CPU). Returns None if unavailable."""
    global _ocr_engine
    if _ocr_engine is None:
        with _ocr_lock:
            if _ocr_engine is None:
                try:
                    from paddleocr import PaddleOCR
                    _ocr_engine = PaddleOCR(use_angle_cls=True, lang='ch', show_log=False)
                except Exception as e:
                    logger.warning(f"[OCR] PaddleOCR unavailable: {e}")
                    _ocr_engine = False  # sentinel: don't retry every call
    return _ocr_engine if _ocr_engine else None

def _ocr_image(filepath):
    """OCR a single image file, return recognized text (multi-line)."""
    if not KB_OCR_ENABLED:
        return ""
    engine = _get_ocr_engine()
    if engine is None:
        return ""
    try:
        result = engine.ocr(filepath, cls=True)
        lines = []
        if result:
            for page in result:
                if not page:
                    continue
                for item in page:
                    if isinstance(item, (list, tuple)) and len(item) >= 2:
                        txt = item[1][0] if isinstance(item[1], (list, tuple)) else str(item[1])
                        if txt and str(txt).strip():
                            lines.append(str(txt).strip())
        return "\n".join(lines)
    except TypeError:
        # PaddleOCR 3.x API: engine.predict(img) returns OCRResult objects
        try:
            result = engine.predict(filepath)
            lines = []
            for res in result:
                txts = getattr(res, "texts", None) or getattr(res, "rec_texts", None)
                if txts:
                    lines.extend(str(t) for t in txts if str(t).strip())
            return "\n".join(lines)
        except Exception as e:
            logger.warning(f"[OCR] 3.x predict failed on {filepath}: {e}")
            return ""
    except Exception as e:
        logger.warning(f"[OCR] failed on {filepath}: {e}")
        return ""

def _ocr_pdf(filepath):
    """OCR a scanned PDF (no text layer): render each page to image, then OCR."""
    if not KB_OCR_ENABLED:
        return ""
    try:
        import fitz
        doc = fitz.open(filepath)
        parts = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=OCR_DPI)
            tmp = os.path.join(tempfile.gettempdir(), f"_kb_ocr_{os.getpid()}_{i}.png")
            try:
                pix.save(tmp)
                t = _ocr_image(tmp)
                if t.strip():
                    parts.append(f"[Page {i+1}]\n{t}")
            finally:
                try: os.remove(tmp)
                except OSError: pass
        doc.close()
        return "\n\n".join(parts)
    except Exception as e:
        logger.warning(f"[OCR] pdf failed {filepath}: {e}")
        return ""

def _parse_image(filepath):
    return _ocr_image(filepath)

def _parse_pdf(filepath):
    try:
        import fitz
        doc = fitz.open(filepath)
        pages_text = []
        for page in doc:
            t = page.get_text()
            if t and t.strip():
                pages_text.append(t)
        doc.close()
        text = "\n".join(pages_text)
        if text.strip():
            return text
        # No text layer → scanned PDF, fall back to OCR
        logger.info(f"[OCR] {os.path.basename(filepath)} has no text layer, running OCR")
        return _ocr_pdf(filepath)
    except ImportError:
        r = subprocess.run(["pdftotext", "-layout", filepath, "-"], capture_output=True, text=True, timeout=60)
        if r.stdout.strip():
            return r.stdout
        return _ocr_pdf(filepath)

def _parse_docx(filepath):
    try:
        import docx
        return "\n".join(p.text for p in docx.Document(filepath).paragraphs)
    except ImportError:
        r = subprocess.run(["python3", "-c", f"import docx;print('\\n'.join(p.text for p in docx.Document('{filepath}').paragraphs))"],
                           capture_output=True, text=True, timeout=30)
        return r.stdout

def _parse_xlsx(filepath):
    import openpyxl
    wb = openpyxl.load_workbook(filepath, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"[Sheet: {name}]")
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows_data.append(cells)
        if not rows_data: continue
        # Convert table to natural language for better embeddings
        if len(rows_data) >= 2:
            header = rows_data[0]
            parts.append(f"Columns: {' | '.join(header)}")
            for r in rows_data[1:]:
                parts.append(" | ".join(f"{header[i]}: {r[i]}" for i in range(min(len(header), len(r))) if i < len(r) and str(r[i]).strip()))
        else:
            for r in rows_data:
                parts.append(" ".join(str(c) for c in r))
    return "\n".join(parts)

def _parse_csv(filepath):
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows: return ""
    header = rows[0]
    parts = [f"Columns: {' | '.join(header)}"]
    for r in rows[1:]:
        parts.append(" | ".join(f"{header[i]}: {r[i]}" for i in range(min(len(header), len(r)))))
    return "\n".join(parts)

def _parse_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        parts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text_frame.text)
            parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
        return "\n\n".join(parts)
    except ImportError:
        return ""

def _parse_json_file(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, ensure_ascii=False, indent=2)

PARSERS = {
    ".txt": _parse_txt, ".md": _parse_md, ".html": _parse_html_file, ".htm": _parse_html_file,
    ".pdf": _parse_pdf, ".docx": _parse_docx, ".xlsx": _parse_xlsx,
    ".csv": _parse_csv, ".pptx": _parse_pptx, ".json": _parse_json_file,
    ".png": _parse_image, ".jpg": _parse_image, ".jpeg": _parse_image,
    ".bmp": _parse_image, ".webp": _parse_image, ".tiff": _parse_image,
}
ALLOWED_EXT = set(PARSERS.keys())

# ── Document Management ─────────────────────────────────
def add_document(title, filepath, ext, dept="general"):
    """Add a new document. If title already exists, create as a new version."""
    import uuid as _uuid
    lang = 'en' if filepath.lower().endswith('-en.md') else 'zh'
    text = PARSERS.get(ext, _parse_txt)(filepath)
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()
    # Check if same-title doc exists for versioning
    existing = db_query("SELECT id, version_group, version FROM documents WHERE title=? ORDER BY created_at DESC LIMIT 1", (title,))
    if existing:
        version_group = existing[0]["version_group"]
        new_version = (existing[0]["version"] or 0) + 1
    else:
        version_group = str(_uuid.uuid4())
        new_version = 1
    cur.execute("INSERT INTO documents (title, filename, filepath, filetype, department, version_group, version) VALUES (?,?,?,?,?,?,?)",
            (title, os.path.basename(filepath), filepath, ext, dept, version_group, new_version))
    conn.commit()
    doc_id = cur.lastrowid
    conn.close()
    # Record in doc_versions
    db_exec("INSERT INTO doc_versions (group_id, doc_id, version) VALUES (?,?,?)",
            (version_group, doc_id, new_version))
    chunks = smart_chunk(text)
    if not chunks:
        logger.error(f"[add_document] Empty text after parsing {filepath}, removing ghost document doc_id={doc_id}")
        db_exec("DELETE FROM documents WHERE id=?", (doc_id,))
        return doc_id, 0
    for i, ch in enumerate(chunks):
        emb = ollama_embed(ch)
        tagged_ch = f"[{title}] {ch}"
        db_exec("INSERT INTO chunks (doc_id, chunk_index, text, embedding, lang) VALUES (?,?,?,?,?)",
                (doc_id, i, tagged_ch, json.dumps(emb), lang))
        bm25.add(doc_id, i, tagged_ch)
    return doc_id, len(chunks)

def delete_document(doc_id):
    bm25.remove(doc_id)
    rows = db_query("SELECT filepath FROM documents WHERE id=?", (doc_id,))
    db_exec("DELETE FROM documents WHERE id=?", (doc_id,))
    if rows and os.path.isfile(rows[0]["filepath"]):
        try: os.remove(rows[0]["filepath"])
        except OSError: pass
    cache_invalidate()

def scan_nas():
    imported = skipped = errors = removed = 0
    existing = set()
    for r in db_query("SELECT filepath FROM documents"):
        existing.add(os.path.normpath(r["filepath"]))
    disk_files = set()
    for root, dirs, files in os.walk(NAS_DOCS_DIR):
        for fname in files:
            ext = os.path.splitext(fname)[1].lower()
            if ext not in ALLOWED_EXT: continue
            fp = os.path.normpath(os.path.join(root, fname))
            disk_files.add(fp)
            if fp in existing: skipped += 1; continue
            title = os.path.splitext(fname)[0]
            dept = os.path.basename(root) if root != NAS_DOCS_DIR else "general"
            try:
                add_document(title, fp, ext, dept)
                imported += 1
            except Exception as e:
                logger.error(f"[Scan] FAILED {fp}: {e}")
                errors += 1
    for db_path in list(existing):
        if db_path not in disk_files:
            rows = db_query("SELECT id FROM documents WHERE filepath=?", (db_path,))
            for r in rows:
                delete_document(r["id"])
                removed += 1
    return {"imported": imported, "skipped": skipped, "errors": errors, "removed": removed}

# ── Chat helpers ─────────────────────────────────────────
def add_message(sid, role, content):
    db_exec("INSERT INTO sessions (session_id, role, content) VALUES (?,?,?)", (sid, role, content))

def get_recent_messages(sid, limit=8):
    rows = db_query("SELECT role, content FROM sessions WHERE session_id=? ORDER BY id DESC LIMIT ?", (sid, limit))
    return list(reversed(rows))

# ── OpenTelemetry (Jaeger) ──────────────────────────────
# DISABLED: 4318 端口无服务，span 导出持续失败刷日志，且可能阻塞请求
# from opentelemetry import trace
# from opentelemetry.sdk.trace import TracerProvider
# from opentelemetry.sdk.trace.export import BatchSpanProcessor
# from opentelemetry.exporter.otlp.proto.http.trace_exporter import OTLPSpanExporter
# from opentelemetry.instrumentation.fastapi import FastAPIInstrumentor
#
#
# _provider = TracerProvider()
# _provider.add_span_processor(BatchSpanProcessor(OTLPSpanExporter(endpoint=JAEGER_ENDPOINT)))
# trace.set_tracer_provider(_provider)

# ── FastAPI App ─────────────────────────────────────────
from fastapi import FastAPI, Request, HTTPException, Query, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse, HTMLResponse
import uvicorn

app = FastAPI(title="MINISFORUM KB V10")

# FastAPIInstrumentor.instrument_app(app)  # DISABLED: 见上方 OpenTelemetry 注释

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Background Periodic Scan ────────────────────────────
@app.on_event("startup")
async def startup_periodic_scan():
    asyncio.create_task(_periodic_scan())

async def _periodic_scan(interval: int = 60):
    """每隔 interval 秒扫描一次 NAS 目录，自动导入新文件并重建索引"""
    while True:
        await asyncio.sleep(interval)
        try:
            result = scan_nas()
            if result.get("imported", 0) > 0 or result.get("removed", 0) > 0:
                print(f"[V10] Auto-scan: +{result['imported']} -{result['removed']}, rebuilding indexes...", flush=True)
                bm25.rebuild()
                _build_embedding_index()
        except Exception as e:
            print(f"[V10] Auto-scan error: {e}", flush=True)

# ── Auth Helpers ────────────────────────────────────────
def decode_jwt(token: str):
    try:
        from jose import jwt as _jwt
        payload = _jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALG])
        return payload.get("sub")
    except Exception:
        return None

def get_user(request: Request):
    auth = request.headers.get("Authorization", "")
    if not auth.startswith("Bearer "):
        return None
    username = decode_jwt(auth[7:])
    if not username:
        return None
    rows = db_query("SELECT * FROM users WHERE username=?", (username,))
    return rows[0] if rows else None

def require_user(request: Request):
    user = get_user(request)
    if not user:
        raise HTTPException(status_code=401, detail="请先登录")
    return user

# ── RAG Pipeline (shared) ───────────────────────────────
def build_rag_context(question: str, user: dict, top_k: int = TOP_K, lang=None):
    """Returns (context_text, results, sources, cache_key) or raises HTTPException."""
    # Cache key
    cache_key = hashlib.md5(f"{question}|{user.get('department','')}|{user.get('role','')}|{lang}".encode()).hexdigest()

    t_start = time.time()

    # Language detection
    is_chinese = bool(re.search(r'[一-鿿]', question))
    search_query = question
    chinese_chars = len(re.findall(r'[\u4e00-\u9fff]', question))
    _product_codes = [
        re.sub(r'[-_\s]', '', t).lower()
        for t in re.findall(r'[A-Za-z0-9]{2,}', question)
        if not re.search(r'[\u4e00-\u9fff]', t) and len(re.sub(r'[-_\s]', '', t)) >= 2
        and not re.sub(r'[-_\s]', '', t).isdigit()
    ]
    is_product_query = bool(_product_codes)
    _non_product_kws = ["公司介绍", "公司简介", "公司背景", "公司历史", "公司概况", "企业介绍",
                        "company intro", "about the company", "company overview"]
    _intro_kws = ["介绍", "简介", "about", "introduction", "overview"]
    _product_kws = ["产品", "参数", "规格", "配置", "型号", "spec", "config", "product"]
    _company_kw_match = any(kw in question.lower() for kw in _non_product_kws)
    _intro_kw_match = any(kw in question.lower() for kw in _intro_kws) and not any(kw in question.lower() for kw in _product_kws)
    is_company_query = _company_kw_match or _intro_kw_match
    if is_company_query:
        _product_codes = []
    # Skip translation if: has English brand/product names, mixed Chinese-English, or question is short
    has_english = bool(re.search(r'[A-Za-z]{2,}', question))
    is_short = len(question) < 20
    # [Optimization] Translation step disabled — skips extra LLM call, search_query stays as original question
    # if chinese_chars >= 4 and not _product_codes and not has_english and not is_short:
    #     try:
    #         translated = ollama_chat([
    #             {"role": "system", "content": "Translate the following Chinese question into English search keywords for document retrieval. Return only the English keywords, no explanation."},
    #             {"role": "user", "content": question}
    #         ], max_tokens=80, temperature=0.0)
    #         if translated and translated.strip():
    #             search_query = translated.strip()
    #     except Exception: pass

    t_after_trans = time.time()
    logger.info("[Perf] translation/lang-detect: %.2fs", t_after_trans - t_start)

    # When query contains specific product codes, boost BM25 (exact keyword match) over vector (semantic)
    # to prevent generic product spec similarity from drowning out precise model matching
    bm25_weight = 2.5 if _product_codes else 1.0
    search_dept = None if (is_company_query or _product_codes) else user.get("department")
    logger.info("[Diag] product_codes=%s search_dept=%s bm25_weight=%.1f query=%s", _product_codes, search_dept, bm25_weight, search_query)
    # Diag: dump matching document chunks for product code queries
    if _product_codes:
        try:
            for pc in _product_codes:
                pc_docs = db_query(
                    "SELECT id,title FROM documents WHERE title LIKE ?",
                    (f"%{pc}%",))
                for d in pc_docs:
                    chks = db_query("SELECT chunk_index,text FROM chunks WHERE doc_id=? ORDER BY chunk_index", (d["id"],))
                    logger.info("[Diag] doc_id=%s title=%s chunks=%d", d["id"], d["title"], len(chks))
                    for c in chks:
                        logger.info("[Diag]   chk[%d]: %s", c["chunk_index"], c["text"][:150])
        except Exception as diag_err:
            logger.warning("[Diag] dump failed (non-fatal): %s", diag_err)
    results = hybrid_search(search_query, top_k, search_dept, user.get("role"), bm25_weight, lang)

    t_after_search = time.time()
    logger.info("[Perf] hybrid_search: %.2fs", t_after_search - t_after_trans)

    # Non-product query -> exclude product spec document chunks
    if not is_product_query and results:
        all_doc_ids = list({k[0] for k, _ in results})
        product_doc_ids = set()
        if all_doc_ids:
            placeholders = ','.join('?' for _ in all_doc_ids)
            doc_rows = db_query(f"SELECT id, title, filepath FROM documents WHERE id IN ({placeholders})", all_doc_ids)
            for row in doc_rows:
                title = row.get("title", "")
                fp = row.get("filepath", "")
                if '_Product_' in fp or '_product_' in fp.lower():
                    product_doc_ids.add(row["id"])
                elif re.search(r'[A-Z0-9]{3,}', title):
                    product_doc_ids.add(row["id"])
        if product_doc_ids:
            non_product = [(k, v) for k, v in results if k[0] not in product_doc_ids]
            if non_product:
                dropped = len(results) - len(non_product)
                results = non_product
                logger.info("[Intent] Non-product query: excluded %d product-doc chunks, %d remaining", dropped, len(results))
            else:
                logger.info("[Intent] Non-product query: all %d chunks are product docs, keeping original results", len(results))

    # Department filter
    all_depts = [r["department"] for r in db_query("SELECT DISTINCT department FROM documents WHERE department != ''")]
    detected_dept = None
    q_lower = question.lower()
    if any(kw in q_lower for kw in _non_product_kws):
        _product_codes = []
        detected_dept = None  # 不做部门限制，全库搜索
    elif any(kw in q_lower for kw in ["产品", "规格", "参数", "配置", "型号", "product", "spec"]):
        detected_dept = next((d for d in all_depts if "产品" in d), None)
    elif any(kw in q_lower for kw in ["技术", "架构", "方案", "白皮书", "technical", "architecture"]):
        detected_dept = next((d for d in all_depts if any(t in d for t in ["技术", "研发"])), None)
    elif _product_codes:
        detected_dept = next((d for d in all_depts if "产品" in d), None)
    if detected_dept:
        filtered = []
        for k, v in results:
            doc_row = db_query("SELECT department FROM documents WHERE id=?", (k[0],))
            if doc_row:
                doc_dept = doc_row[0]["department"]
                doc_dept_base = doc_dept.replace("_ZH","").replace("_EN","").replace("_zh","").replace("_en","")
                dept_base = detected_dept.replace("_ZH","").replace("_EN","")
                if doc_dept_base == dept_base:
                    filtered.append((k, v))
        if filtered: results = filtered

    # Document-level dedup: keep top 2 most relevant documents to avoid cross-contamination
    if results:
        doc_scores = {}
        for key, info in results:
            doc_id = key[0]
            if doc_id not in doc_scores:
                doc_scores[doc_id] = 0.0
            doc_scores[doc_id] += info.get("rrf_score", 0)
        top_docs = sorted(doc_scores.items(), key=lambda x: -x[1])[:2]
        top_doc_ids = {d for d, _ in top_docs}
        results = [(k, v) for k, v in results if k[0] in top_doc_ids]

    # Language filter: when lang is specified, prefer docs matching the requested language
    if lang and results:
        lang_suffix = "-EN" if lang == "en" else "-CN"
        filtered = [(k, v) for k, v in results if get_document_title(k[0]).upper().endswith(lang_suffix)]
        if filtered:
            results = filtered
            logger.info("[LangFilter] kept %d docs matching lang=%s (suffix=%s)", len(filtered), lang, lang_suffix)

    # Fuzzy fallback
    fuzzy_contexts = []
    has_vec_score = any(info.get("vec_score") is not None for _, info in results)
    max_sim = max((info.get("vec_score") or 0 for _, info in results), default=0) if has_vec_score else None

    # Preamble injection — 最多追加 2 个额外 chunk，且结果已 >= 10 条时跳过
    if results and len(results) < 10:
        top_doc_id = results[0][0][0]
        seen_keys = {key for key, _ in results}
        for row in db_query("SELECT chunk_index, text FROM chunks WHERE doc_id=? ORDER BY chunk_index ASC LIMIT 2", (top_doc_id,)):
            key = (top_doc_id, row["chunk_index"])
            if key not in seen_keys:
                results.append((key, {"text": row["text"], "vec_score": None, "rrf_score": 0.0, "bm25_rank": -1}))

    # Fast fuzzy fallback: extract keywords from query and search BM25 only (no LLM blocking)
    if not results or (max_sim is not None and max_sim < 0.12):
        keywords = re.findall(r'[\u4e00-\u9fff]{2,}|[a-zA-Z0-9]{2,}', question)
        seen_fuzzy = set()
        for kw in keywords[:6]:
            for doc_id, chunk_id, text, bm_score in bm25.search(kw, top_k=2):
                if user.get("role") not in ("admin", "superadmin"):
                    di_rows = db_query("SELECT department FROM documents WHERE id=?", (doc_id,))
                    if di_rows and user.get("department") and di_rows[0].get("department") != user.get("department"):
                        continue
                key = (doc_id, chunk_id)
                if key not in seen_fuzzy:
                    seen_fuzzy.add(key)
                    fuzzy_contexts.append({"key": key, "title": get_document_title(doc_id), "text": text})

    # Apply language filter to fuzzy contexts as well
    if lang and fuzzy_contexts:
        lang_suffix = "-EN" if lang == "en" else "-CN"
        filtered_fc = [fc for fc in fuzzy_contexts if fc["title"].upper().endswith(lang_suffix)]
        if filtered_fc:
            fuzzy_contexts = filtered_fc
            logger.info("[LangFilter] fuzzy: kept %d docs matching lang=%s", len(filtered_fc), lang)

    if not results and not fuzzy_contexts:
        raise HTTPException(status_code=200, detail="No documents found matching your question.")

    # Product code filter (with fuzzy I↔1, O↔0 matching and empty-result fallback)
    _no_exact_match = False
    if _product_codes:
        def _norm_title(key):
            return re.sub(r'[-_\s]', '', get_document_title(key[0])).lower()

        def _code_variants(code):
            """Generate fuzzy variants: I↔1, O↔0 substitutions for OCR/typo tolerance."""
            variants = {code}
            if '1' in code or 'i' in code:
                variants.add(code.replace('1', 'i').replace('0', 'o'))
                # Also try letter→digit
                variants.add(code.replace('i', '1').replace('o', '0'))
            if 'o' in code or '0' in code:
                variants.add(code.replace('o', '0'))
                variants.add(code.replace('0', 'o'))
            return variants

        def _doc_matches(key):
            title = _norm_title(key)
            return all(
                any(v in title for v in _code_variants(code))
                for code in _product_codes
            )

        matched = [(k, v) for k, v in results if _doc_matches(k)]
        _no_exact_match = bool(_product_codes and not matched)
        if matched:
            results = matched
        if fuzzy_contexts:
            matched_fc = [
                fc for fc in fuzzy_contexts
                if all(
                    any(v in re.sub(r'[-_\s]', '', fc['title']).lower() for v in _code_variants(code))
                    for code in _product_codes
                )
            ]
            if matched_fc:
                fuzzy_contexts = matched_fc

    # Build context text
    contexts = []
    for i, (key, info) in enumerate(results):
        contexts.append(f"[Doc {i+1}: {get_document_title(key[0])}]\n{info['text']}")
    context_text = "\n\n---\n\n".join(contexts)

    if fuzzy_contexts:
        if context_text: context_text += "\n\n---\n\n"
        for j, fc in enumerate(fuzzy_contexts):
            context_text += f"[Fuzzy Doc {j+1}: {fc['title']}]\n{fc['text']}\n\n---\n\n"

    context_text = context_text.rstrip("- \n")
    if len(context_text) > MAX_CONTEXT_CHARS:
        context_text = context_text[:MAX_CONTEXT_CHARS]
        # 尝试在最后一个完整段落处截断
        last_sep = context_text.rfind("\n\n---")
        if last_sep > MAX_CONTEXT_CHARS * 0.7:
            context_text = context_text[:last_sep]
    if not context_text:
        context_text = "未找到相关文档。"
    if _no_exact_match:
        context_text = (f"[系统提示] 知识库中未找到关于 {'、'.join(_product_codes)} 的精确文档，以下为可能相关内容（非指定型号）。\n\n" + context_text)

    t_after_ctx = time.time()
    logger.info("[Perf] context_build (dept_filter+fuzzy+preamble+truncate): %.2fs", t_after_ctx - t_after_search)

    # Build sources
    sources = []
    seen_src = set()
    for key, info in results:
        doc_id = key[0]
        if doc_id not in seen_src:
            seen_src.add(doc_id)
            sources.append({"doc_id": doc_id, "title": get_document_title(doc_id), "score": info.get("rrf_score", info.get("vec_score", 0) or 0)})
    for fc in fuzzy_contexts:
        doc_id = fc["key"][0]
        if doc_id not in seen_src:
            seen_src.add(doc_id)
            sources.append({"doc_id": doc_id, "title": fc["title"], "score": 0.01})

    logger.info("[Perf] build_rag_context total: %.2fs", time.time() - t_start)
    return context_text, sources, cache_key

# ── API Routes ──────────────────────────────────────────
@app.post("/api/login")
async def login(request: Request):
    data = await request.json()
    u = data.get("username", "").strip()
    p = data.get("password", "")
    rows = db_query("SELECT * FROM users WHERE username=?", (u,))
    if not rows:
        raise HTTPException(status_code=401, detail="用户名或密码错误")
    user = rows[0]
    if user["status"] != "approved":
        raise HTTPException(status_code=403, detail="账号未审核通过")
    try:
        from passlib.hash import bcrypt as _bcrypt
        if not _bcrypt.verify(p, user["password_hash"]):
            raise HTTPException(status_code=401, detail="用户名或密码错误")
    except Exception:
        if p != user["password_hash"]:
            raise HTTPException(status_code=401, detail="用户名或密码错误")
    from jose import jwt as _jwt
    token = _jwt.encode({"sub": u, "role": user["role"], "dept": user["department"]}, JWT_SECRET, algorithm=JWT_ALG)
    return {"ok": True, "token": token, "user": {"username": u, "department": user["department"], "role": user["role"]}}

@app.post("/api/register")
async def register(request: Request):
    data = await request.json()
    u = data.get("username", "").strip()
    p = data.get("password", "")
    d = data.get("department", "general")
    if not u or len(p) < 6:
        raise HTTPException(status_code=400, detail="用户名不能为空，密码至少6位")
    if db_query("SELECT id FROM users WHERE username=?", (u,)):
        raise HTTPException(status_code=409, detail="用户名已存在")
    try:
        from passlib.hash import bcrypt as _bcrypt
        pwd_hash = _bcrypt.hash(p)
    except ImportError:
        pwd_hash = p
    db_exec("INSERT INTO users (username, password_hash, department) VALUES (?,?,?)", (u, pwd_hash, d))
    return {"ok": True, "message": "注册成功，请等待管理员审核"}

@app.get("/api/health")
async def health():
    stats = {}
    for tbl in ("documents", "chunks", "users", "sessions"):
        r = db_query(f"SELECT COUNT(*) as cnt FROM {tbl}")
        stats[tbl] = r[0]["cnt"] if r else 0
    return {"ok": True, "version": "v10", "stats": stats,
            "models": {"ask": KB_MODEL_ASK, "ask_en": KB_MODEL_ASK_EN, "chat": KB_MODEL_CHAT,
                       "compare": KB_MODEL_COMPARE, "extractor": KB_MODEL_EXTRACTOR,
                       "embed": EMBED_MODEL}}

@app.get("/api/docs")
async def list_docs(request: Request, page: int = 1, page_size: int = 100):
    require_user(request)
    # Show only latest version of each doc by default; pass ?all=1 to show all versions
    show_all = request.query_params.get("all") == "1"
    if show_all:
        rows = db_query("SELECT id, title, filename, department, created_at, version_group, version FROM documents ORDER BY id DESC LIMIT ? OFFSET ?",
                        (page_size, (page - 1) * page_size))
    else:
        rows = db_query("""
            SELECT d.id, d.title, d.filename, d.department, d.created_at, d.version_group, d.version
            FROM documents d
            INNER JOIN (
                SELECT version_group, MAX(version) as max_v
                FROM documents
                GROUP BY version_group
            ) latest ON d.version_group = latest.version_group AND d.version = latest.max_v
            ORDER BY d.id DESC LIMIT ? OFFSET ?
        """, (page_size, (page - 1) * page_size))
    return {"ok": True, "docs": rows}

@app.delete("/api/docs/{doc_id}")
async def delete_doc(request: Request, doc_id: int):
    require_user(request)
    delete_document(doc_id)
    cache_invalidate()
    return {"ok": True}

@app.get("/api/docs/{doc_id}")
async def get_doc(request: Request, doc_id: int):
    require_user(request)
    rows = db_query("SELECT id, title, filename, filepath, department, created_at FROM documents WHERE id=?", (doc_id,))
    if not rows:
        return {"ok": False, "error": "Document not found"}
    doc = rows[0]
    chunks = db_query("SELECT text FROM chunks WHERE doc_id=? ORDER BY chunk_index", (doc_id,))
    doc["content"] = "\n\n".join(c["text"] for c in chunks) if chunks else ""
    return {"ok": True, "doc": doc}

@app.get("/api/docs/{doc_id}/file")
async def serve_doc_file(request: Request, doc_id: int):
    """Serve the original document file for download/viewing."""
    from fastapi.responses import FileResponse
    require_user(request)
    rows = db_query("SELECT filepath, filename FROM documents WHERE id=?", (doc_id,))
    if not rows:
        return {"ok": False, "error": "Document not found"}
    fp = rows[0]["filepath"]
    fn = rows[0]["filename"]
    if not fp or not os.path.isfile(fp):
        return {"ok": False, "error": "File not found on disk"}
    return FileResponse(fp, filename=fn)

# ── Document Versions ───────────────────────────────────
@app.get("/api/docs/{doc_id}/versions")
async def list_versions(request: Request, doc_id: int):
    """List all versions of a document by version_group."""
    require_user(request)
    row = db_query("SELECT version_group FROM documents WHERE id=?", (doc_id,))
    if not row:
        return {"ok": False, "error": "Document not found"}
    group_id = row[0]["version_group"]
    versions = db_query(
        "SELECT d.id, d.title, d.filename, d.version, d.created_at "
        "FROM documents d WHERE d.version_group=? ORDER BY d.version ASC",
        (group_id,))
    return {"ok": True, "versions": versions, "group_id": group_id}

@app.post("/api/docs/versions/diff")
async def version_diff(request: Request):
    """Diff between two document versions by doc_id."""
    require_user(request)
    data = await request.json()
    v1_id = data.get("v1")
    v2_id = data.get("v2")
    if not v1_id or not v2_id:
        return {"ok": False, "error": "v1 and v2 (doc_id) are required"}
    rows1 = db_query("SELECT id, title, version, created_at FROM documents WHERE id=?", (v1_id,))
    rows2 = db_query("SELECT id, title, version, created_at FROM documents WHERE id=?", (v2_id,))
    if not rows1 or not rows2:
        return {"ok": False, "error": "One or both documents not found"}
    
    chunks1 = db_query("SELECT text FROM chunks WHERE doc_id=? ORDER BY chunk_index", (v1_id,))
    chunks2 = db_query("SELECT text FROM chunks WHERE doc_id=? ORDER BY chunk_index", (v2_id,))
    old_lines = "\n".join(c["text"] for c in (chunks1 or [])).splitlines()
    new_lines = "\n".join(c["text"] for c in (chunks2 or [])).splitlines()
    
    import difflib
    diff_result = []
    for line in difflib.unified_diff(old_lines, new_lines,
                                      fromfile=f"{rows1[0]['title']} v{rows1[0]['version']}",
                                      tofile=f"{rows2[0]['title']} v{rows2[0]['version']}",
                                      lineterm=""):
        diff_result.append(line)
    
    added = sum(1 for l in diff_result if l.startswith("+") and not l.startswith("+++"))
    removed = sum(1 for l in diff_result if l.startswith("-") and not l.startswith("---"))
    return {"ok": True,
            "v1": {"id": v1_id, "title": rows1[0]["title"], "version": rows1[0]["version"]},
            "v2": {"id": v2_id, "title": rows2[0]["title"], "version": rows2[0]["version"]},
            "added": added, "removed": removed,
            "diff": "\n".join(diff_result[:200])}

@app.get("/api/kb/scan")
async def kb_scan(request: Request):
    require_user(request)
    result = scan_nas()
    bm25.rebuild()
    cache_invalidate()
    return {"ok": True, **result}

@app.get("/api/kb/nuke")
async def kb_nuke(request: Request):
    """Delete ALL documents/chunks, rebuild from scratch. No auth for recovery."""
    rows = db_query("SELECT id FROM documents")
    for r in rows:
        bm25.remove(r["id"])
    db_exec("DELETE FROM chunks")
    db_exec("DELETE FROM documents")
    db_exec("DELETE FROM queries")
    cache_invalidate()
    result = scan_nas()
    bm25.rebuild()
    return {"ok": True, "purged": len(rows), **result}

@app.get("/api/kb/dept-list")
async def list_all_depts():
    """调试用：返回所有文档的 department 和 title"""
    rows = db_query("SELECT id, title, department, filetype FROM documents ORDER BY id")
    return {"count": len(rows), "items": [dict(r) for r in rows]}

@app.get("/api/kb/cleanup")
async def kb_cleanup(request: Request):
    """删除孤儿 chunk（chunks 表中 doc_id 在 documents 表中不存在的记录）并重建 BM25"""
    require_user(request)
    orphan_rows = db_query("""
        SELECT DISTINCT c.doc_id FROM chunks c
        LEFT JOIN documents d ON c.doc_id = d.id
        WHERE d.id IS NULL
    """)
    orphan_ids = [r["doc_id"] for r in orphan_rows]
    if not orphan_ids:
        return {"ok": True, "cleaned": 0, "message": "No orphan chunks found"}

    for oid in orphan_ids:
        db_exec("DELETE FROM chunks WHERE doc_id=?", (oid,))
        try:
            bm25.remove(oid)
        except Exception:
            pass

    cache_invalidate()
    logger.info("[Cleanup] removed orphan chunks for doc_ids: %s", orphan_ids)
    return {"ok": True, "cleaned": len(orphan_ids), "orphan_ids": orphan_ids}

# ── KB Ask (RAG) with SSE streaming ────────────────────
@app.post("/api/kb/ask")
@app.post("/api/ai/sms/ask")
async def kb_ask(request: Request):
    user = require_user(request)
    _incr_query()
    data = await request.json()
    question = data.get("question", "").strip()
    session_id = data.get("session_id", "default")
    top_k = data.get("top_k", TOP_K)
    lang = data.get("lang", "zh")

    if not question:
        raise HTTPException(status_code=400, detail="question 不能为空")

    db_exec("INSERT INTO queries (query, session_id) VALUES (?,?)", (question, session_id))

    # Check cache
    cached = cache_get(hashlib.md5(f"{question}|{user.get('department','')}|{user.get('role','')}|{lang}".encode()).hexdigest())
    if cached:
        print(f"[V10 cache hit] {question[:50]}", flush=True)
        def generate_cached():
            yield f"data: {json.dumps({'type': 'sources', 'sources': cached['sources']})}\n\n"
            for chunk in [cached["answer"][i:i+100] for i in range(0, len(cached["answer"]), 100)]:
                yield f"data: {json.dumps({'type': 'token', 'content': chunk})}\n\n"
            yield f"data: {json.dumps({'type': 'done'})}\n\n"
        return StreamingResponse(generate_cached(), media_type="text/event-stream",
                                 headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

    # Build RAG context
    try:
        context_text, sources, cache_key = build_rag_context(question, user, top_k, lang)
    except HTTPException as e:
        if e.status_code == 200:
            fallback_answer = e.detail
            def generate_fallback():
                yield f"data: {json.dumps({'type': 'sources', 'sources': []})}\n\n"
                yield f"data: {json.dumps({'type': 'token', 'content': fallback_answer})}\n\n"
                yield f"data: {json.dumps({'type': 'done'})}\n\n"
            return StreamingResponse(generate_fallback(), media_type="text/event-stream",
                                     headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})
        raise

    if not context_text.strip():
        is_vague = classify_question(question)
        fallback_answer = (
            f"知识库中暂未找到关于「{question}」的明确信息。你的提问比较简短，"
            f"能否补充一下你想了解的具体方面？比如：详细参数、功能特点、使用场景等，我来帮你精准查找。"
        ) if is_vague else "未在知识库中找到相关信息。"
        def generate_fallback():
            yield f"data: {json.dumps({'type': 'sources', 'sources': sources})}\n\n"
            yield f"data: {json.dumps({'type': 'token', 'content': fallback_answer})}\n\n"
            yield f"data: {json.dumps({'type': 'done'})}\n\n"
        return StreamingResponse(generate_fallback(), media_type="text/event-stream",
                                 headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

    # Relevance guard: check if any query term appears in source titles OR chunk content
    query_terms = set(re.findall(r'[a-zA-Z0-9]{2,}|[\u4e00-\u9fff]+', question.lower()))
    if query_terms:
        title_text = " ".join(s["title"] for s in sources).lower()
        combined_text = title_text + " " + context_text.lower()
        combined_terms = set(re.findall(r'[a-zA-Z0-9]+|[\u4e00-\u9fff]+', combined_text))
        if not query_terms & combined_terms:
            logger.info("[Relevance] no query-term overlap with source titles, returning not-found")
            not_found_answer = (
                "No matching documents found in the knowledge base. Please check the product name or try a different query."
                if lang == "en"
                else f"知识库中未找到与「{question}」相关的文档，请检查产品名称或尝试其他关键词。"
            )
            def generate_not_found():
                yield f"data: {json.dumps({'type': 'sources', 'sources': sources})}\n\n"
                yield f"data: {json.dumps({'type': 'token', 'content': not_found_answer})}\n\n"
                yield f"data: {json.dumps({'type': 'done'})}\n\n"
            return StreamingResponse(generate_not_found(), media_type="text/event-stream",
                                     headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

    # Weak results guard: ≤2 sources with low similarity → suggest language switch
    if len(sources) <= 2 and all(s.get("score", 0) < 0.12 for s in sources):
        if lang == "en":
            weak_answer = "No matching documents found in the knowledge base for this language. Please switch to another language and try again."
        else:
            weak_answer = "该问题在当前语言下的知识库中暂无匹配文档，请切换语言后重试。"
        def generate_weak():
            yield f"data: {json.dumps({'type': 'sources', 'sources': sources})}\n\n"
            yield f"data: {json.dumps({'type': 'token', 'content': weak_answer})}\n\n"
            yield f"data: {json.dumps({'type': 'done'})}\n\n"
        return StreamingResponse(generate_weak(), media_type="text/event-stream",
                                 headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

    # Build prompt
    extra_detail = ""
    if lang == "en":
        extra_detail = (
            "CRITICAL: Extract and present EVERY specification, parameter, and technical detail "
            "from the documents below. Do NOT summarize or condense. Match the level of detail "
            "you would provide in a Chinese answer — English answers must be equally thorough. "
            "List ALL numeric values, dimensions, ports, protocols, and features explicitly.\n"
        )
    if lang == "en":
        # UI language is English: ALWAYS answer in English, even if the question is typed in Chinese.
        prompt = (
            "Answer in English. The user's question may be written in Chinese, but you MUST respond in English. "
            "Base your answer ONLY on the documents below. Be comprehensive and well-structured. "
            + extra_detail +
            "Use Markdown tables for specs, parameters, configurations, or comparisons. "
            "Do NOT include a Sources or References section.\n"
            f"Documents:\n{context_text}\n\n"
            f"User question: {question}"
        )
    else:
        prompt = (
            "Answer in Chinese. Base your answer ONLY on the documents below. Be comprehensive and well-structured. "
            "Use Markdown tables for specs, parameters, configurations, or comparisons. "
            "Do NOT include a Sources or References section.\n"
            f"Documents:\n{context_text}\n\n"
            f"User question: {question}"
        )

    # SSE streaming generator (sync → FastAPI auto-runs in thread pool, avoids blocking uvicorn event loop)
    def generate():
        global _query_queue_counter

        # ── Concurrency control: queue position + semaphore ──
        with _query_counter_lock:
            _query_queue_counter += 1
            pos = _query_queue_counter

        if pos > 1:
            yield f"data: {json.dumps({'type': 'queue', 'i18n_key': 'queue_waiting', 'n': pos - 1, 'position': pos})}\n\n"

        _query_semaphore.acquire()
        try:
            with _query_counter_lock:
                _query_queue_counter -= 1
                pos = _query_queue_counter

            if pos > 0:
                yield f"data: {json.dumps({'type': 'queue', 'i18n_key': 'queue_your_turn'})}\n\n"

            t_gen_start = time.time()
            full_answer = ""
            prompt_len = len(prompt)
            logger.info("[KB] Query streaming start | prompt=%d chars | %d sources", prompt_len, len(sources))
            yield f"data: {json.dumps({'type': 'sources', 'sources': sources})}\n\n"
            t_first_token = time.time()
            try:
                for delta_type, delta in ollama_chat_stream(
                    [{"role": "system", "content": prompt}, {"role": "user", "content": question}],
                    model=KB_MODEL_ASK_EN if lang == "en" else KB_MODEL_ASK, temperature=0.2, use_gpu=True
                ):
                    if delta_type == 'token':
                        if not full_answer:
                            t_first = (time.time() - t_first_token) * 1000
                            logger.info("[Perf] FIRST TOKEN after %.0fms | prompt=%d chars (%d sources) | context_build=%.2fs | LLM_prefill=%.2fs",
                                        t_first, prompt_len, len(sources),
                                        t_first_token - t_gen_start, time.time() - t_gen_start)
                        full_answer += delta
                    yield f"data: {json.dumps({'type': delta_type, 'content': delta})}\n\n"
                logger.info("[KB] Query done | answer=%d chars | %.1fs total", len(full_answer), time.time() - t_gen_start)
                if not full_answer:
                    logger.warning("[KB] EMPTY ANSWER — check kb_stream log for reason")
                yield f"data: {json.dumps({'type': 'done'})}\n\n"
                # Cache the result (skip hallucinated refusal answers)
                _refusal_patterns = ["您是否能提供更多", "could you provide more", "请提供更多",
                                     "I don't have enough context", "我没有足够"]
                if not any(p in full_answer for p in _refusal_patterns):
                    cache_set(cache_key, {"answer": full_answer, "sources": sources})
                add_message(session_id, "user", question)
                add_message(session_id, "assistant", full_answer)
            except Exception as e:
                import traceback
                with open(os.path.join(tempfile.gettempdir(), "kb_v10_error_trace.txt"), "a") as f:
                    f.write(f"\n=== KB/ASK ERROR ===\n{traceback.format_exc()}\n")
                yield f"data: {json.dumps({'type': 'error', 'content': traceback.format_exc()[:500]})}\n\n"
        finally:
            _query_semaphore.release()

    return StreamingResponse(generate(), media_type="text/event-stream",
                             headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

# ── Chat with SSE streaming ─────────────────────────────
@app.post("/api/chat")
async def chat(request: Request):
    user = require_user(request)
    data = await request.json()
    msg = data.get("message", "").strip()
    sid = data.get("session_id", "default")
    if not msg:
        raise HTTPException(status_code=400, detail="message 不能为空")
    add_message(sid, "user", msg)
    recent = get_recent_messages(sid, limit=8)
    msgs = [{"role": r["role"], "content": r["content"]} for r in recent]

    def generate():
        global _query_queue_counter

        # ── Concurrency control: queue position + semaphore ──
        with _query_counter_lock:
            _query_queue_counter += 1
            pos = _query_queue_counter

        if pos > 1:
            yield f"data: {json.dumps({'type': 'queue', 'i18n_key': 'queue_waiting', 'n': pos - 1, 'position': pos})}\n\n"

        _query_semaphore.acquire()
        try:
            with _query_counter_lock:
                _query_queue_counter -= 1
                pos = _query_queue_counter

            if pos > 0:
                yield f"data: {json.dumps({'type': 'queue', 'i18n_key': 'queue_your_turn'})}\n\n"

            full_answer = ""
            try:
                for delta_type, delta in ollama_chat_stream(msgs, model=KB_MODEL_CHAT, use_gpu=True):
                    if delta_type == 'token':
                        full_answer += delta
                    yield f"data: {json.dumps({'type': delta_type, 'content': delta})}\n\n"
                logger.info("[Chat] done | answer=%d chars", len(full_answer))
                yield f"data: {json.dumps({'type': 'done'})}\n\n"
                add_message(sid, "assistant", full_answer)
            except Exception as e:
                import traceback
                with open(os.path.join(tempfile.gettempdir(), "kb_v10_error_trace.txt"), "a") as f:
                    f.write(f"\n=== CHAT ERROR ===\n{traceback.format_exc()}\n")
                yield f"data: {json.dumps({'type': 'error', 'content': traceback.format_exc()[:500]})}\n\n"
        finally:
            _query_semaphore.release()

    return StreamingResponse(generate(), media_type="text/event-stream",
                             headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

# ── Compare ─────────────────────────────────────────────
@app.post("/api/kb/compare")
async def compare(request: Request):
    user = require_user(request)
    data = await request.json()
    question = data.get("question", "").strip()
    dimension = data.get("dimension", question)
    doc_ids = data.get("doc_ids", [])
    if not question:
        raise HTTPException(status_code=400, detail="question 不能为空")
    if len(doc_ids) < 2:
        raise HTTPException(status_code=400, detail="至少选择2个文档")

    # Collect contexts per document
    doc_contexts = {}
    for did in doc_ids:
        chunks = db_query("SELECT chunk_index, text FROM chunks WHERE doc_id=? ORDER BY chunk_index", (did,))
        doc_contexts[did] = "\n\n".join(c["text"] for c in chunks)

    context_text = ""
    for did, text in doc_contexts.items():
        title = get_document_title(did)
        context_text += f"\n\n### [{title}]\n{text[:5000]}\n"

    # Single-stage comparison: 跳过 extract_facts（省一次 LLM 调用，避免慢/失败），直接用原文
    facts = context_text[:8000]

    prompt = (
        f"你是企业文档对比分析专家。以下是各文档关于「{dimension}」的提取事实。\n"
        f"请据此生成对比分析：\n"
        f"1. 先给出综合对比表格（Markdown表格）\n"
        f"2. 逐篇分析各自特点\n"
        f"3. 总结和建议\n\n"
        f"提取事实：\n{facts}\n\n"
        f"要求：基于提取的事实，不要编造。若某文档无该维度数据，表格中标注「无数据」。输出务必精简：表格每格一句话，逐篇分析每篇不超过3句，总结不超过5句。"
    )
    try:
        if KB_MODEL_COMPARE == MINISLLM_MODEL:
            answer = minisllm_chat([{"role": "user", "content": prompt}], max_tokens=4096, temperature=0.3)
        else:
            answer = ollama_chat([{"role": "user", "content": prompt}], model=KB_MODEL_COMPARE, max_tokens=1500, base_url=OLLAMA_HOST_COMPARE)
    except Exception as e:
        answer = f"LLM 调用失败: {e}"

    return {"ok": True, "comparison": answer, "dimension": dimension}

@app.post("/api/kb/compare_stream")
async def compare_stream(request: Request):
    user = require_user(request)
    data = await request.json()
    question = data.get("question", "").strip()
    dimension = data.get("dimension", question)
    doc_ids = data.get("doc_ids", [])
    if not question:
        raise HTTPException(status_code=400, detail="question 不能为空")
    if len(doc_ids) < 2:
        raise HTTPException(status_code=400, detail="至少选择2个文档")

    doc_contexts = {}
    for did in doc_ids:
        chunks = db_query("SELECT chunk_index, text FROM chunks WHERE doc_id=? ORDER BY chunk_index", (did,))
        doc_contexts[did] = "\n\n".join(c["text"] for c in chunks)

    context_text = ""
    for did, text in doc_contexts.items():
        title = get_document_title(did)
        context_text += f"\n\n### [{title}]\n{text[:5000]}\n"

    facts = context_text[:8000]

    prompt = (
        f"你是企业文档对比分析专家。以下是各文档关于「{dimension}」的提取事实。\n"
        f"请据此生成对比分析：\n"
        f"1. 先给出综合对比表格（Markdown表格）\n"
        f"2. 逐篇分析各自特点\n"
        f"3. 总结和建议\n\n"
        f"提取事实：\n{facts}\n\n"
        f"要求：基于提取的事实，不要编造。若某文档无该维度数据，表格中标注「无数据」。输出务必精简：表格每格一句话，逐篇分析每篇不超过3句，总结不超过5句。"
    )

    def gen():
        try:
            for kind, text in ollama_chat_stream([{"role": "user", "content": prompt}], model=KB_MODEL_COMPARE, max_tokens=1500, temperature=0.3):
                if kind == "token":
                    yield "data: " + json.dumps({"type": "token", "content": text}, ensure_ascii=False) + "\n\n"
        except Exception as e:
            yield "data: " + json.dumps({"type": "error", "message": str(e)}, ensure_ascii=False) + "\n\n"
        yield "data: " + json.dumps({"type": "done"}, ensure_ascii=False) + "\n\n"

    return StreamingResponse(gen(), media_type="text/event-stream",
                             headers={"X-Accel-Buffering": "no", "Cache-Control": "no-cache"})

# ── Gaps ────────────────────────────────────────────────
@app.get("/api/kb/gaps")
async def gaps(request: Request):
    require_user(request)
    unanswered = db_query("""
        SELECT q.query, COUNT(*) as freq FROM queries q LEFT JOIN sessions s ON q.session_id=s.session_id
        WHERE q.answered=0 GROUP BY q.query ORDER BY freq DESC LIMIT 20
    """)
    gaps_data = [{"query": r["query"], "frequency": r["freq"]} for r in unanswered]
    docs = db_query("SELECT id, title FROM documents")
    doc_list = [f"{d['id']}: {d['title']}" for d in docs[:10]]
    suggestion = "建议补充以下领域的文档以覆盖知识缺口。" if gaps_data else "暂无明显知识缺口。"
    return {"ok": True, "gaps": gaps_data, "suggestion": suggestion, "documents": doc_list}

# ── Feedback ────────────────────────────────────────────
@app.post("/api/kb/feedback")
async def feedback(request: Request):
    require_user(request)
    data = await request.json()
    qid = data.get("query_id")
    rating = data.get("rating", "")
    if qid:
        db_exec("INSERT INTO feedback (query_id, rating) VALUES (?,?)", (qid, rating))
    return {"ok": True}

# ── Backup ──────────────────────────────────────────────
@app.post("/api/kb/backup")
async def backup(request: Request):
    user = require_user(request)
    if user.get("role") not in ("superadmin", "admin"):
        raise HTTPException(status_code=403, detail="仅管理员可执行备份")
    try:
        proc = subprocess.run(["sh", "/kb_persist/backup.sh"], capture_output=True, text=True, timeout=300)
        output = proc.stdout + proc.stderr
        if proc.returncode != 0:
            raise HTTPException(status_code=500, detail=f"备份脚本执行失败: {output}")
        filename = "kb_backup_unknown.tar.gz"
        for line in output.splitlines():
            if line.strip().endswith(".tar.gz"):
                filename = os.path.basename(line.strip())
                break
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=500, detail="备份脚本超时（300秒）")
    except FileNotFoundError:
        raise HTTPException(status_code=500, detail="备份脚本不存在: /kb_persist/backup.sh")
    return {"success": True, "message": "备份完成", "filename": filename}

# ── User Management ─────────────────────────────────────
@app.get("/api/users")
async def list_users(request: Request):
    user = require_user(request)
    if user.get("role") not in ("admin", "superadmin"):
        raise HTTPException(status_code=403, detail="权限不足")
    rows = db_query("SELECT id, username, department, role, status FROM users")
    return {"ok": True, "users": rows}

@app.post("/api/users/approve")
async def user_action(request: Request):
    user = require_user(request)
    if user.get("role") not in ("admin", "superadmin"):
        raise HTTPException(status_code=403, detail="权限不足")
    data = await request.json()
    uid = data.get("user_id")
    action = data.get("action")
    if action == "approve":
        db_exec("UPDATE users SET status='approved' WHERE id=?", (uid,))
    elif action == "reject":
        db_exec("DELETE FROM users WHERE id=?", (uid,))
    return {"ok": True}

# ── Graph (Entity search) ───────────────────────────────
@app.get("/api/kb/graph")
async def graph(request: Request, q: str = ""):
    require_user(request)
    if not q:
        return {"ok": True, "results": []}
    results, seen_entities = [], set()
    chunks = db_query("SELECT c.text, d.title FROM chunks c JOIN documents d ON c.doc_id=d.id LIMIT 300")
    for chunk in chunks:
        tokens = re.findall(r'[A-Z][A-Z0-9_-]{2,}|[A-Z][a-z]+(?:[ -][A-Z][a-z]+)+', chunk["text"][:2000])
        entities = [t for t in tokens if t.lower() != q.lower() and t not in seen_entities]
        related = []
        for ent in entities[:8]:
            seen_entities.add(ent)
            related.append({"type": "mentions", "target": ent, "source": chunk["title"]})
        if related:
            results.append({"entity": q, "related": related})
            break
    if not results:
        results.append({"entity": q, "related": []})
    return {"ok": True, "results": results}

# ── Stale documents ─────────────────────────────────────
@app.get("/api/kb/stale")
async def stale(request: Request):
    require_user(request)
    rows = db_query("SELECT id, title, filename, department, created_at FROM documents ORDER BY created_at ASC")
    stale_list = []
    cutoff = time.time() - 90 * 86400
    for r in rows:
        ts = r.get("created_at") or ""
        try:
            t = time.mktime(time.strptime(ts[:19], "%Y-%m-%d %H:%M:%S"))
            if t < cutoff:
                r["updated_at"] = ts[:19]
                stale_list.append(r)
        except: continue
    return {"ok": True, "docs": stale_list[:50]}

# ── Conflicts ───────────────────────────────────────────
@app.get("/api/kb/conflicts")
async def conflicts(request: Request):
    require_user(request)
    docs = db_query("SELECT id, title, department FROM documents LIMIT 50")
    doc_texts = {}
    for doc in docs:
        chunks = db_query("SELECT text FROM chunks WHERE doc_id=? LIMIT 10", (doc["id"],))
        doc_texts[doc["id"]] = (doc["title"], " ".join(c["text"] or "" for c in chunks)[:5000])
    conflicts_list = []
    field_patterns = {
        "TDP": r'(?:TDP|功耗)[：:\s]*(\d+[-~]\d+|\d+)\s*W',
        "RAM": r'(?:RAM|内存|Memory)[：:\s]*(\d+)\s*GB',
        "Cores": r'(?:Core|核心|Cores?)[：:\s]*(\d+)[Cores\s]*',
        "Price": r'(?:价格|Price|售价)[：:\s]*\$?(\d+)',
        "Storage": r'(?:Storage|存储|SSD|HDD)[：:\s]*(\d+)\s*(GB|TB)',
    }
    for field, pattern in field_patterns.items():
        values = {}
        for did, (title, text) in doc_texts.items():
            m = re.search(pattern, text, re.IGNORECASE)
            if m:
                val = m.group(1).strip()
                if val not in values: values[val] = []
                values[val].append(title)
        if len(values) >= 2:
            val_list = [{"value": k, "doc": ", ".join(v)} for k, v in values.items()]
            conflicts_list.append({"field": field, "values": val_list})
    return {"ok": True, "conflicts": conflicts_list[:20]}

# ── Upload ──────────────────────────────────────────────
# /api/kb/upload is deprecated; use /api/kb/upload-multipart instead
@app.post("/api/kb/upload-multipart")
async def upload_multipart(request: Request):
    user = require_user(request)
    content_type = request.headers.get("content-type", "")
    if "multipart/form-data" not in content_type:
        raise HTTPException(status_code=400, detail="需要 multipart/form-data")

    form = await request.form()
    dept = form.get("department", "general") if hasattr(form, "get") else "general"
    uploaded_files = form.getlist("file") if hasattr(form, "getlist") else [v for k, v in form.items() if hasattr(v, "filename")]

    imported = 0
    for f in uploaded_files:
        if not hasattr(f, "filename"): continue
        ext = os.path.splitext(f.filename)[1].lower()
        if ext not in ALLOWED_EXT: continue
        title = os.path.splitext(f.filename)[0]
        save_path = os.path.join(NAS_DOCS_DIR, f.filename)
        with open(save_path, "wb") as out:
            out.write(await f.read())
        try:
            add_document(title, save_path, ext, str(dept))
            imported += 1
        except Exception: pass

    bm25.rebuild()
    cache_invalidate()
    return {"ok": True, "imported": imported}

# ── History ────────────────────────────────────────────
@app.get("/api/history")
async def history(request: Request, limit: int = 50):
    require_user(request)
    rows = db_query(
        "SELECT s.id, s.session_id, s.role, s.content, s.created_at "
        "FROM sessions s WHERE s.role IN ('user','assistant') "
        "ORDER BY s.id DESC LIMIT ?", (limit*2,))
    sessions_map = {}
    for row in reversed(rows):
        sid = row["session_id"]
        if sid not in sessions_map:
            sessions_map[sid] = {"session_id": sid, "pairs": [], "last_at": row["created_at"]}
        sessions_map[sid]["pairs"].append({"role": row["role"], "content": row["content"][:300]})
    history_list = []
    for s in sorted(sessions_map.values(), key=lambda x: x["last_at"] or "", reverse=True):
        pairs = s["pairs"]
        i = 0
        while i < len(pairs):
            if pairs[i]["role"] == "user" and i+1 < len(pairs) and pairs[i+1]["role"] == "assistant":
                history_list.append({
                    "question": pairs[i]["content"],
                    "answer": pairs[i+1]["content"],
                    "session_id": s["session_id"],
                    "created_at": s["last_at"]
                })
                i += 2
            else:
                i += 1
    return {"ok": True, "history": history_list[:limit]}

# ── Keyword Search ─────────────────────────────────────
@app.get("/api/search")
async def keyword_search(request: Request, q: str = "", limit: int = 30, lang: str = "zh"):
    require_user(request)
    if not q.strip():
        return {"ok": True, "results": []}
    terms = q.strip().split()
    conditions = " AND ".join([
        "(c.text LIKE ? OR REPLACE(REPLACE(REPLACE(c.text, '-', ''), '_', ''), ' ', '') LIKE ?)"
        for _ in terms
    ])
    params = []
    for t in terms:
        params.append("%"+t+"%")
        params.append("%"+re.sub(r'[-_\s]', '', t)+"%")
    params.append(limit)
    rows = db_query(
        f"SELECT c.text, d.title as doc_title, d.filename, d.id as doc_id "
        f"FROM chunks c JOIN documents d ON c.doc_id=d.id "
        f"WHERE c.lang=? AND {conditions} ORDER BY c.id DESC LIMIT ?",
        [lang] + params)
    results = []
    for row in rows:
        text = row["text"]
        snippet = text[:400]
        for t in terms:
            idx = text.lower().find(t.lower())
            if idx >= 0:
                start = max(0, idx-60)
                snippet = ("..." if start>0 else "") + text[start:idx+len(t)+120] + ("..." if idx+len(t)+120<len(text) else "")
                break
            # fallback: separator-insensitive match ("MS01" ↔ "MS-01")
            sep_pattern = r'[-_\s]*'.join(re.escape(c) for c in t)
            m = re.search(sep_pattern, text, re.IGNORECASE)
            if m:
                start = max(0, m.start()-60)
                snippet = ("..." if start>0 else "") + text[start:m.end()+120] + ("..." if m.end()+120<len(text) else "")
                break
        for t in terms:
            snippet = snippet.replace(t, f"<mark>{t}</mark>")
            snippet = snippet.replace(t.capitalize(), f"<mark>{t.capitalize()}</mark>")
        results.append({
            "doc_title": row["doc_title"],
            "filename": row["filename"],
            "doc_id": row["doc_id"],
            "snippet": snippet
        })
    return {"ok": True, "results": results, "query": q}

# ── Related Questions ──────────────────────────────────
@app.get("/api/related-questions")
async def related_questions(request: Request):
    require_user(request)
    recent = db_query(
        "SELECT content FROM sessions WHERE role='user' ORDER BY id DESC LIMIT 5")
    if not recent:
        return {"ok": True, "questions": []}
    recent_qs = [r["content"][:100] for r in recent]
    # Build suggestions from doc titles + recent queries
    docs = db_query("SELECT title FROM documents ORDER BY id DESC LIMIT 10")
    titles = [d["title"] for d in docs]
    suggestions = []
    for t in titles:
        if len(suggestions) >= 6:
            break
        suggestions.append(f"{t} 的主要参数是什么？")
    # Pad remaining slots from defaults
    defaults = [
        "KB 里一共有多少文档？",
        "最近导入了哪些新文档？",
        "当前知识库覆盖了哪些主题？",
        "哪些文档最后被修改过？",
        "帮我总结 KB 的核心内容",
        "搜索所有包含 'spec' 的文档",
    ]
    for q in defaults:
        if len(suggestions) >= 6:
            break
        if q not in suggestions:
            suggestions.append(q)
    return {"ok": True, "questions": suggestions[:6]}

# ── Doc Diff ───────────────────────────────────────────
@app.post("/api/docs/diff")
async def doc_diff(request: Request):
    require_user(request)
    form = await request.form()
    uploaded_files = form.getlist("files") if hasattr(form, "getlist") else [v for k, v in form.items() if hasattr(v, "filename")]
    if len(uploaded_files) != 2:
        return {"ok": False, "error": "需要上传恰好两个文件（旧版和新版）"}
    contents = []
    for f in uploaded_files:
        raw = await f.read()
        text = raw.decode("utf-8", errors="ignore")[:200_000]
        contents.append(text)
    # Simple line diff
    old_lines = contents[0].splitlines()
    new_lines = contents[1].splitlines()
    import difflib
    diff_result = []
    for line in difflib.unified_diff(old_lines, new_lines,
                                      fromfile=uploaded_files[0].filename,
                                      tofile=uploaded_files[1].filename,
                                      lineterm=""):
        diff_result.append(line)
    # Count stats
    added = sum(1 for l in diff_result if l.startswith("+") and not l.startswith("+++"))
    removed = sum(1 for l in diff_result if l.startswith("-") and not l.startswith("---"))
    return {"ok": True, "added": added, "removed": removed,
            "old_file": uploaded_files[0].filename, "new_file": uploaded_files[1].filename,
            "diff": "\n".join(diff_result[:200])}

# ── Spec Compare ──────────────────────────────────────

def parse_specmd(text):
    """Extract flat key-value pairs from document text.
    Handles patterns: key: value, **key**: value, key | value, - key: value
    Returns list of {key: str, value: str}"""
    # Remove YAML front matter (AIGC metadata block: Label, ContentProducer, etc.)
    text = re.sub(r'^---\s*\n.*?---\s*\n', '', text, flags=re.DOTALL)
    items = []
    for line in text.splitlines():
        line = line.strip()
        if not line or len(line) < 3:
            continue
        # Skip pure markdown headers (not parameter lines)
        if line.startswith('#') and ':' not in line and '|' not in line:
            continue
        # Pattern 1: **key** | value or **key**: value
        m = re.match(r'\*\*(.+?)\*\*\s*[:|]\s*(.+)', line)
        if m:
            key = m.group(1).strip()
            val = m.group(2).strip()
            items.append({"key": key, "value": val})
            continue
        # Pattern 2: key: value
        m = re.match(r'^([^:|#\n]{2,50}?)\s*:\s*(.+)', line)
        if m:
            key = m.group(1).strip()
            val = m.group(2).strip()
            if not key.startswith('!') and not key.startswith('['):
                items.append({"key": key, "value": val})
            continue
        # Pattern 3: key | value (table row, first two columns)
        m = re.match(r'^\|?\s*([^|]+?)\s*\|\s*([^|]+)', line)
        if m:
            key = m.group(1).strip()
            val = m.group(2).strip()
            if key and not key.startswith('-') and not key.startswith('#'):
                items.append({"key": key, "value": val})
    return items


def generate_spec_comparison(file_results):
    """file_results: [{"filename": str, "items": [{"key": str, "value": str}]}]
    Returns: {"rows": [{"key": str, "values": [str]}], "products": [str], "product_count": int}"""
    # Collect all unique keys, preserving order
    all_keys = []
    seen = set()
    for fr in file_results:
        for item in fr["items"]:
            if item["key"] not in seen:
                all_keys.append(item["key"])
                seen.add(item["key"])
    
    rows = []
    for key in all_keys:
        values = []
        for fr in file_results:
            val = "—"
            for item in fr["items"]:
                if item["key"] == key:
                    val = item["value"]
                    break
            values.append(val)
        rows.append({"key": key, "values": values})
    
    products = [fr["filename"] for fr in file_results]
    return {"rows": rows, "products": products, "product_count": len(file_results)}


@app.post("/api/spec/compare")
async def spec_compare(request: Request):
    require_user(request)
    try:
        form = await request.form()
        uploaded_files = form.getlist("files") if hasattr(form, "getlist") else [v for k, v in form.items() if hasattr(v, "filename")]
        if len(uploaded_files) < 2 or len(uploaded_files) > 10:
            return {"ok": False, "error": "需要上传 2-10 个文件"}
        file_data = []
        for f in uploaded_files:
            raw = await f.read()
            text = raw.decode("utf-8", errors="ignore")[:200_000]
            file_data.append({"filename": f.filename, "text": text})
        # Try SpecMD parsing
        product_results = []
        all_valid = True
        for fd in file_data:
            items = parse_specmd(fd["text"])
            product_results.append({"filename": fd["filename"], "items": items, "text": fd["text"]})
            if not items:
                all_valid = False
        mode = "spec"
        if not all_valid:
            mode = "llm"
            for fd in file_data:
                try:
                    prompt = (
                        "Extract all parameter key-value pairs from the following product specification document.\n"
                        "Return ONLY a valid JSON array (no markdown fences, no extra text):\n"
                        '[{"key":"参数名","value":"值"},...]\n\n'
                        "Include every parameter you find: technical specs, dimensions, features, etc.\n"
                        "If a value spans multiple short values, combine them into one string.\n\n"
                        "Document text:\n" + fd["text"][:8000] + "\n\nJSON:"
                    )
                    resp = ollama_chat([{"role": "user", "content": prompt}], model=KB_MODEL_EXTRACTOR, max_tokens=2048)
                    json_match = re.search(r'\[[\s\S]*\]', resp)
                    items = []
                    if json_match:
                        raw_json = json_match.group()
                        # Strip markdown code fences if present
                        raw_json = re.sub(r'^```(?:json)?\s*\n?', '', raw_json, flags=re.MULTILINE)
                        raw_json = re.sub(r'\n?```\s*$', '', raw_json, flags=re.MULTILINE)
                        raw_json = raw_json.strip()
                        try:
                            items = json.loads(raw_json)
                            if not isinstance(items, list):
                                items = []
                        except json.JSONDecodeError:
                            logger.warning(f"JSON parse failed for {fd['filename']}")
                    product_results.append({"filename": fd["filename"], "items": items, "text": fd["text"]})
                except Exception as e:
                    logger.warning(f"LLM extraction failed for {fd['filename']}: {e}")
                    product_results.append({"filename": fd["filename"], "items": [], "text": fd["text"]})
        comparison = generate_spec_comparison(product_results)
        analysis = ""
        try:
            summary_prompt = (
                "Compare the following products and give a brief analysis in Chinese (≤200 characters). "
                "Highlight the most notable differences.\n\n"
                + json.dumps([{"name": fr["filename"],
                               "items": fr["items"][:20]} for fr in product_results], ensure_ascii=False)
                + "\n\nBrief analysis:"
            )
            if KB_MODEL_COMPARE == MINISLLM_MODEL:
                analysis = minisllm_chat([{"role": "user", "content": summary_prompt}], max_tokens=300, temperature=0.3)
            else:
                analysis = ollama_chat([{"role": "user", "content": summary_prompt}], model=KB_MODEL_COMPARE, max_tokens=300, base_url=OLLAMA_HOST_COMPARE)
            analysis = analysis.strip()[:300]
        except Exception as e:
            logger.warning(f"LLM summary failed: {e}")
        return {"ok": True, "mode": mode, "products": comparison["products"],
                "table": comparison, "analysis": analysis}
    except Exception as e:
        logger.error(f"spec_compare crashed: {e}", exc_info=True)
        return {"ok": False, "error": f"Spec compare failed: {str(e)[:200]}"}

# ── Docs Export ────────────────────────────────────────
@app.get("/api/docs/export")
async def docs_export(request: Request):
    require_user(request)
    docs = db_query("SELECT * FROM documents ORDER BY id")
    result = []
    for d in docs:
        chunks = db_query("SELECT text FROM chunks WHERE doc_id=? ORDER BY id", (d["id"],))
        result.append({
            "id": d["id"], "title": d.get("title",""), "filename": d.get("filename",""),
            "department": d.get("department",""), "created_at": d.get("created_at",""),
            "chunks": [c["text"] for c in chunks]
        })
    return {"ok": True, "total_docs": len(result), "docs": result}

# ── Stats ───────────────────────────────────────────────
@app.get("/api/stats")
async def stats(request: Request):
    require_user(request)
    docs_count = db_query("SELECT COUNT(*) as cnt FROM documents")[0]["cnt"]
    chunks_count = db_query("SELECT COUNT(*) as cnt FROM chunks")[0]["cnt"]
    users_count = db_query("SELECT COUNT(*) as cnt FROM users")[0]["cnt"]
    with _stats_lock:
        q_today = _stats["queries_today"]
    cache_hits = 0
    with _cache_lock:
        cache_hits = len(_cache)
    return {"ok": True, "docs": docs_count, "chunks": chunks_count, "users": users_count,
            "queries_today": q_today, "cache_entries": cache_hits}

# ── Frontend ────────────────────────────────────────────
@app.get("/")
async def serve_frontend():
    # Inline the same V9 frontend HTML (slightly updated for V10 streaming)
    return HTMLResponse(FRONTEND_HTML)

FRONTEND_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>MINISFORUM KB</title>
<style>
:root {
  --bg: #08080f; --glass-bg: rgba(255,255,255,0.03); --glass-border: rgba(255,255,255,0.06);
  --surface: rgba(255,255,255,0.04); --text: #e4e4ed; --text-dim: #8b8b9e; --text-muted: #5c5c72;
  --primary: #6366f1; --primary-end: #8b5cf6; --accent: #6366f1; --warm: #f59e0b; --success: #10b981; --danger: #ef4444; --warning: #f59e0b;
  --radius: 14px; --radius-lg: 20px; --radius-pill: 9999px;
  --shadow: 0 4px 16px rgba(0,0,0,0.5), 0 1px 4px rgba(0,0,0,0.4);
  --shadow-lg: 0 12px 40px rgba(0,0,0,0.6), 0 4px 12px rgba(0,0,0,0.4);
  --transition: 0.2s cubic-bezier(0.4, 0, 0.2, 1);
}
* { margin:0; padding:0; box-sizing:border-box; }
html, body { min-height:100vh; }
body {
  font: 400 14px/1.6 'Inter', 'SF Pro Display', system-ui, -apple-system, sans-serif;
  background: var(--bg); color: var(--text); display:flex; overflow-y:auto;
  letter-spacing: -0.01em; -webkit-font-smoothing: antialiased;
}
::-webkit-scrollbar { width:6px; }
::-webkit-scrollbar-track { background:transparent; }
::-webkit-scrollbar-thumb { background:rgba(255,255,255,0.08); border-radius:var(--radius-pill); }
.sidebar {
  width:280px; background:var(--glass-bg); border-right:1px solid var(--glass-border);
  display:flex; flex-direction:column; flex-shrink:0;
  position:sticky; top:0; height:100vh;
  backdrop-filter:blur(24px) saturate(140%);
}
.sidebar-header {
  padding:20px 20px 16px; border-bottom:1px solid var(--glass-border);
  font-size:18px; font-weight:700; letter-spacing:-0.02em;
  background:linear-gradient(135deg, var(--primary), var(--primary-end));
  -webkit-background-clip:text; -webkit-text-fill-color:transparent; background-clip:text;
}
.sidebar-header .ver { font-size:11px; font-weight:400; -webkit-text-fill-color:var(--text-dim); margin-left:6px; }
.sidebar-user {
  padding:14px 20px; border-bottom:1px solid var(--glass-border);
  font-size:12px; color:var(--text-dim); display:flex; justify-content:space-between; align-items:center;
}
.sidebar-user .logout { color:var(--warm); cursor:pointer; font-weight:500; padding:4px 10px; border-radius:var(--radius-pill); }
.stats-dash { display:grid; grid-template-columns:1fr 1fr; gap:6px; }
.stat-card {
  background:rgba(255,255,255,0.04); border-radius:10px; padding:10px 8px;
  text-align:center; transition:background .2s;
}
.stat-card:hover { background:rgba(255,255,255,0.08); }
.stat-val { display:block; font-size:20px; font-weight:700; color:var(--accent); line-height:1.2; }
.stat-lbl { display:block; font-size:10px; color:var(--text-muted); margin-top:2px; }
.feat-list { padding:4px 16px 8px; display:flex; flex-direction:column; gap:2px; }
.feat-item {
  display:flex; align-items:center; gap:10px; padding:10px 12px;
  border-radius:10px; cursor:pointer; transition:background .15s;
}
.feat-item:hover { background:rgba(255,255,255,0.05); }
.feat-item.active { background:rgba(99,102,241,0.12); }
.feat-icon { font-size:18px; flex-shrink:0; width:24px; text-align:center; }
.feat-name { display:block; font-size:13px; font-weight:600; color:var(--text); line-height:1.3; }
.feat-desc { display:block; font-size:10px; color:var(--text-muted); }
.related-q { padding:4px 16px 12px; border-top:1px solid var(--glass-border); margin-top:6px; }
.related-q-title { font-size:10px; color:var(--text-dim); text-transform:uppercase; letter-spacing:0.5px; margin:10px 0 8px; }
.related-q-tag { display:inline-block; margin:3px 4px; padding:4px 10px; background:rgba(99,102,241,0.08); border-radius:999px; font-size:11px; color:var(--text-dim); cursor:pointer; transition:all .2s; word-break:break-all; }
.related-q-tag:hover { background:rgba(99,102,241,0.18); color:var(--text); }
.nav-tabs { display:flex; border-bottom:1px solid var(--glass-border); padding:10px 14px; }
.nav-tab {
  flex:1; padding:12px 28px;
  background:linear-gradient(135deg, #6366f1, #8b5cf6);
  border:none; border-radius:10px;
  color:#fff; font:inherit; font-size:15px; font-weight:700;
  letter-spacing:0.5px;
  cursor:pointer; text-align:center; transition:all 0.25s cubic-bezier(0.4,0,0.2,1);
  box-shadow:0 2px 16px rgba(99,102,241,0.35), 0 0 24px rgba(139,92,246,0.12);
}
.nav-tab:hover {
  background:linear-gradient(135deg, #7c3aed, #a78bfa);
  box-shadow:0 4px 28px rgba(99,102,241,0.55), 0 0 40px rgba(139,92,246,0.25);
  transform:translateY(-2px) scale(1.03);
}
.nav-tab:active {
  transform:translateY(0) scale(0.98);
  box-shadow:0 1px 8px rgba(99,102,241,0.3);
  transition:all 0.1s ease;
}
.nav-tab.active {
  background:linear-gradient(135deg, #7c3aed, #a78bfa);
  box-shadow:0 3px 24px rgba(124,58,237,0.5), 0 0 48px rgba(139,92,246,0.3);
}
.main {
  flex:1; display:flex; flex-direction:column; min-width:0;
}
.chat-area {
  flex:1; padding:24px 28px; display:flex; flex-direction:column; gap:16px;
}
.message {
  max-width:80%; padding:14px 18px; border-radius:var(--radius-lg);
  font-size:14px; line-height:1.7; word-wrap:break-word;
}
.message.user {
  align-self:flex-end; background:linear-gradient(135deg, rgba(99,102,241,0.25), rgba(139,92,246,0.2));
  border-bottom-right-radius:6px; word-break:break-word; overflow-wrap:break-word; overflow:hidden;
}
.message.assistant {
  align-self:flex-start; background:var(--surface); border-bottom-left-radius:6px;
  word-break:break-word; overflow-wrap:break-word; overflow:hidden;
}
.message.assistant .kb-reasoning {
  padding: 4px 10px; border-left: 2px solid rgba(99,102,241,0.12);
  color: var(--text-muted); font-size: 0.65rem; line-height: 1.2; font-style: italic;
  opacity: 0.4; max-width: 300px; width: 100%; max-height: 34px;
  overflow-y: auto; white-space: pre-wrap; word-break: break-all;
  scrollbar-width: thin; scrollbar-color: rgba(255,255,255,0.06) transparent;
}

.message.assistant table { width:100%; border-collapse:collapse; margin:12px 0; }
.message.assistant th, .message.assistant td { padding:8px 12px; border:1px solid var(--glass-border); text-align:left; }
.message.assistant th { background:rgba(99,102,241,0.15); }
.input-area {
  padding:16px 28px; border-top:1px solid var(--glass-border); display:flex; gap:10px;
  background:var(--glass-bg); backdrop-filter:blur(24px);
  position:sticky; bottom:0; z-index:10;
}
.input-area input {
  flex:1; padding:12px 18px; border-radius:var(--radius-pill); border:1px solid var(--glass-border);
  background:var(--surface); color:var(--text); font:inherit; outline:none;
  transition:border-color var(--transition);
}
.input-area input:focus { border-color:var(--primary); }
.input-area button {
  padding:12px 24px; border-radius:var(--radius-pill); border:none;
  background:linear-gradient(135deg, var(--primary), var(--primary-end));
  color:#fff; font:inherit; font-weight:600; cursor:pointer;
  transition:all var(--transition);
}
.input-area button:hover { box-shadow:0 4px 20px rgba(99,102,241,0.4); transform:translateY(-1px); }
.input-area input:disabled { background:rgba(255,255,255,0.02); opacity:0.5; cursor:not-allowed; }
.input-area button:disabled { opacity:0.6; cursor:not-allowed; transform:none !important; }
.input-area button.loading {
  animation: btnPulse 1.5s ease-in-out infinite;
  box-shadow: 0 0 20px rgba(99,102,241,0.5);
}
@keyframes btnPulse {
  0%, 100% { opacity: 0.8; }
  50% { opacity: 1; }
}
/* ── 进度卡片（紧凑版，纯进度条）── */
.kb-progress-card {
  align-self: flex-start; max-width: 300px; width:100%;
  padding: 10px 16px;
  background: var(--surface); border: 1px solid var(--glass-border);
  border-radius: var(--radius); margin-top: 4px;
}
.kb-progress-bar-wrap {
  height: 6px; background: rgba(255,255,255,0.05);
  border-radius: 999px; overflow: hidden; margin-bottom: 8px;
}
.kb-progress-fill {
  height: 100%; width: 0%;
  background: linear-gradient(90deg, var(--primary), var(--primary-end));
  border-radius: 999px;
  transition: width 0.8s ease;
  box-shadow: 0 0 12px rgba(99,102,241,0.35), 0 0 3px rgba(99,102,241,0.15);
}
/* 加载中呼吸辉光 */
.kb-progress-fill.loading {
  animation: kbProgPulse 2s ease-in-out infinite;
}
@keyframes kbProgPulse {
  0%, 100% {
    box-shadow: 0 0 8px rgba(99,102,241,0.25), 0 0 2px rgba(99,102,241,0.1);
    opacity: 0.7;
  }
  50% {
    box-shadow: 0 0 20px rgba(99,102,241,0.55), 0 0 8px rgba(139,92,246,0.3);
    opacity: 1;
  }
}
.kb-progress-meta {
  display: flex; align-items: center;
}
.kb-progress-status {
  font-size: 11px; color: var(--text-dim);
}
@keyframes kbProgFadeOut {
  to { opacity: 0; transform: translateY(-4px); }
}
.kb-progress-card.done {
  animation: kbProgFadeOut 0.3s 0.1s ease forwards;
  pointer-events: none;
}

.toast {
  position:fixed; top:20px; right:20px; padding:12px 24px; border-radius:var(--radius);
  background:var(--surface); backdrop-filter:blur(20px); border:1px solid var(--glass-border);
  color:var(--text); font-size:13px; z-index:9999; animation:fadeIn 0.3s ease;
  box-shadow:var(--shadow-lg);
}
@keyframes fadeIn { from { opacity:0; transform:translateY(-10px); } to { opacity:1; transform:translateY(0); } }
.login-overlay {
  position:fixed; inset:0; background:rgba(0,0,0,0.6); display:flex; align-items:center; justify-content:center;
  z-index:999; backdrop-filter:blur(8px);
}
.login-box {
  background:var(--glass-bg); border:1px solid var(--glass-border); border-radius:var(--radius-lg);
  padding:32px; width:360px; box-shadow:var(--shadow-lg); backdrop-filter:blur(24px);
}
.login-box h2 { margin-bottom:20px; font-size:20px; font-weight:700; letter-spacing:-0.02em; }
.login-box input {
  width:100%; padding:12px; margin-bottom:12px; border-radius:var(--radius);
  border:1px solid var(--glass-border); background:var(--surface); color:var(--text); font:inherit;
  outline:none;
}
.login-box input:focus { border-color:var(--primary); }
.login-box button {
  width:100%; padding:12px; border-radius:var(--radius); border:none;
  background:linear-gradient(135deg, var(--primary), var(--primary-end));
  color:#fff; font:inherit; font-weight:600; cursor:pointer;
}
.sources-bar {
  padding:8px 28px; border-top:1px solid var(--glass-border); font-size:11px; color:var(--text-muted);
  display:flex; gap:12px; flex-wrap:wrap;
}
.sources-bar span { background:rgba(99,102,241,0.1); padding:2px 8px; border-radius:var(--radius-pill); }
.source-chip { cursor:pointer; transition:all 0.15s; padding:6px 14px !important; border:1px solid rgba(99,102,241,0.3); display:inline-flex; align-items:center; gap:8px; font-size:12px; border-radius:var(--radius); background:rgba(99,102,241,0.08); }
.source-chip .source-chip-label { background:transparent; padding:0; color:var(--text); font-weight:500; max-width:240px; overflow:hidden; text-overflow:ellipsis; white-space:nowrap; }
.source-chip .source-chip-score { background:rgba(99,102,241,0.15); padding:1px 6px; border-radius:var(--radius-pill); color:var(--primary); font-size:10px; }
.source-chip:hover { background:rgba(99,102,241,0.18) !important; border-color:rgba(99,102,241,0.6); }

@keyframes spin { to { transform:rotate(360deg); } }

select {
  -webkit-appearance:none; appearance:none;
  background-image:url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' width='12' height='12' viewBox='0 0 12 12'%3E%3Cpath fill='%238b8b9e' d='M6 8L1 3h10z'/%3E%3C/svg%3E");
  background-repeat:no-repeat; background-position:right 10px center;
  padding-right:30px !important; cursor:pointer;
}

/* ── Diff 展示优化：分组折叠 + 长行换行 + 仅看变更 ── */
.diff-toolbar{position:sticky;top:0;z-index:5;display:flex;align-items:center;gap:8px;flex-wrap:wrap;padding:8px 12px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius);margin-bottom:10px;font-size:12px}
.diff-toolbar .btn{padding:4px 10px;background:transparent;border:1px solid var(--glass-border);border-radius:var(--radius-pill);color:var(--text-dim);font-size:11px;cursor:pointer}
.diff-toolbar .btn:hover{color:var(--text);border-color:var(--primary)}
.diff-hunk{border:1px solid var(--glass-border);border-radius:var(--radius);margin-bottom:8px;overflow:hidden;background:var(--surface)}
.diff-hunk-head{display:flex;align-items:center;gap:8px;padding:6px 10px;cursor:pointer;user-select:none;background:rgba(255,255,255,0.02);border-bottom:1px solid var(--glass-border)}
.diff-hunk-head:hover{background:rgba(255,255,255,0.05)}
.diff-hunk-arrow{font-size:10px;color:var(--text-dim);transition:transform .15s;display:inline-block}
.diff-hunk.collapsed .diff-hunk-arrow{transform:rotate(-90deg)}
.diff-hunk-pos{font-family:monospace;font-size:11px;color:var(--text-dim);flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.diff-hunk-badge{font-family:monospace;font-size:11px;font-weight:600}
.diff-hunk-body{padding:4px 0}
.diff-row{font-family:monospace;font-size:11px;line-height:1.5;padding:1px 10px;white-space:pre-wrap;word-break:break-all;overflow-wrap:break-word}
.diff-row:hover{background:rgba(255,255,255,0.04)}
.diff-add{background:rgba(16,185,129,0.10);color:#10b981}
.diff-del{background:rgba(239,68,68,0.10);color:#ef4444}
.diff-ctx{color:var(--text-muted)}
.diff-hunk.collapsed .diff-hunk-body{display:none}
.diff-hide-ctx .diff-ctx{display:none}
.diff-pair{display:flex;align-items:flex-start;gap:8px;padding:2px 10px}
.diff-pair-old{flex:1;color:#ef4444;text-decoration:line-through;opacity:.85;background:rgba(239,68,68,0.08);border-radius:4px;padding:0 6px;white-space:pre-wrap;word-break:break-all;overflow-wrap:break-word}
.diff-pair-new{flex:1;color:#10b981;background:rgba(16,185,129,0.10);border-radius:4px;padding:0 6px;white-space:pre-wrap;word-break:break-all;overflow-wrap:break-word}
.diff-pair-arrow{color:var(--text-dim);flex-shrink:0;line-height:1.5}
.diff-row b{font-weight:600;color:inherit}

</style>
<script src="https://cdn.jsdelivr.net/npm/marked/marked.min.js"></script>
</head>
<body>

<div id="login-modal" class="login-overlay" style="display:none">
  <div class="login-box">
    <h2>MINISFORUM KB</h2>
    <input id="login-username" placeholder="Username" autocomplete="username">
    <input id="login-password" type="password" placeholder="Password" autocomplete="current-password">
    <button onclick="doLogin()">Sign In</button>
    <p style="margin-top:12px;font-size:11px;color:var(--text-muted);text-align:center">Default: admin / admin123</p>
  </div>
</div>

<div class="sidebar" id="sidebar" style="display:none">
  <div class="sidebar-header" style="display:flex;justify-content:space-between;align-items:center">
    <span>MINISFORUM KB</span>
    <select id="lang-switch" onchange="setLang(this.value)" style="background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius-pill);color:var(--text-dim);padding:3px 8px;font-size:11px;cursor:pointer;outline:none">
      <option value="en">EN</option>
      <option value="zh">中文</option>
    </select>
  </div>
  <div class="sidebar-user"><span id="user-info"></span><span class="logout" onclick="doLogout()">Logout</span></div>
  <div class="nav-tabs">
    <button class="nav-tab active" data-page="ask" onclick="switchPage('ask')">Ask</button>
  </div>
  <div class="stats-dash" id="stats-dash">
    <div class="stat-card"><span class="stat-val" id="st-docs">-</span><span class="stat-lbl">Docs</span></div>
    <div class="stat-card"><span class="stat-val" id="st-chunks">-</span><span class="stat-lbl">Chunks</span></div>
    <div class="stat-card"><span class="stat-val" id="st-queries">-</span><span class="stat-lbl">Queries</span></div>
    <div class="stat-card"><span class="stat-val" id="st-users">-</span><span class="stat-lbl">Users</span></div>
  </div>
  <div class="feat-list">
    <div class="feat-item" data-page="history" onclick="switchPage('history')">
      <span class="feat-name" data-i18n="feat_history">历史问答</span>
    </div>
    <div class="feat-item" data-page="docs" onclick="switchPage('docs')">
      <span class="feat-name" data-i18n="feat_docs">文档管理</span>
    </div>
    <div class="feat-item" data-page="search" onclick="switchPage('search')">
      <span class="feat-name" data-i18n="feat_search">关键词搜索</span>
    </div>
    <div class="feat-item" data-page="compare" onclick="switchPage('compare')">
      <span class="feat-name" data-i18n="feat_compare">文档对比</span>
    </div>
  </div>
  <div class="related-q" id="related-q">
    <div class="related-q-title">你可能还想问</div>
    <div id="related-tags"></div>
  </div>
</div>

<div class="main" id="main-panel" style="display:none">
  <div class="chat-area" id="chat-area"></div>
  <div class="sources-bar" id="sources-bar"></div>
  <div class="input-area">
    <input id="ask-input" placeholder="Ask MINISFORUM KB..." onkeydown="if(event.key==='Enter'&&!isComposing)askQuestion()">
    <button onclick="askQuestion()">Send</button>
  </div>
</div>

<script>
// ── Internationalization ──
const _i18n = {
  username:{en:"Username",zh:"用户名"},
  password:{en:"Password",zh:"密码"},
  sign_in:{en:"Sign In",zh:"登录"},
  login_hint:{en:"Default: admin / admin123",zh:"默认: admin / admin123"},
  login_failed:{en:"Login failed",zh:"登录失败"},
  logout:{en:"Logout",zh:"退出登录"},
  ask:{en:"Ask",zh:"问答"},
  history:{en:"History",zh:"历史问答"},
  docs:{en:"Docs",zh:"文档管理"},
  search:{en:"Search",zh:"全文搜索"},
  history_desc:{en:"Q&A History",zh:"历史问答"},
  docs_desc:{en:"Document Management",zh:"文档管理"},
  search_desc:{en:"Full-text Search",zh:"关键词全文搜索"},
  feat_history:{en:"History",zh:"历史问答"},
  feat_docs:{en:"Document Management",zh:"文档管理"},
  feat_search:{en:"Keyword Search",zh:"关键词搜索"},
  feat_compare:{en:"Docs Compare",zh:"文档对比"},
  compare_title:{en:"Document Comparison",zh:"文档对比"},
  compare_desc:{en:"Select documents from knowledge base and specify a dimension to compare",zh:"从知识库选择文档并指定对比维度，生成分析对比表格"},
  compare_sel_docs:{en:"Select documents",zh:"选择文档"},
  compare_dimension:{en:"Comparison dimension (e.g. CPU / price / features)",zh:"对比维度（如 CPU / 价格 / 功能）"},
  compare_btn:{en:"Compare",zh:"开始对比"},
  compare_comparing:{en:"Analyzing...",zh:"分析中..."},
  compare_no_docs:{en:"Select at least 2 documents",zh:"至少选择2个文档"},
  compare_no_dim:{en:"Please enter a comparison dimension",zh:"请输入对比维度"},
  version_label:{en:"v{0}",zh:"v{0}"},
  version_history:{en:"Version History",zh:"版本历史"},
  version_diff_btn:{en:"Diff",zh:"对比差异"},
  version_only_one:{en:"Only 1 version — upload a newer version to enable diff.",zh:"仅有 1 个版本，上传新版本后可对比差异"},
  version_diff_title:{en:"Version Diff",zh:"版本差异对比"},
  version_diff_pick:{en:"Select two versions to compare",zh:"选择两个版本进行对比"},
  docs_back:{en:"Back to list",zh:"返回列表"},
  stats_docs:{en:"Docs",zh:"文档"},
  stats_chunks:{en:"Chunks",zh:"分块"},
  stats_queries:{en:"Queries",zh:"查询"},
  stats_users:{en:"Users",zh:"用户"},
  related_title:{en:"You may also ask",zh:"你可能还想问"},
  input_placeholder:{en:"Ask MINISFORUM KB...",zh:"输入问题..."},
  send:{en:"Send",zh:"发送"},
  step_search:{en:"Searching Docs",zh:"检索匹配文档"},
  step_analyze:{en:"Analyzing Content",zh:"分析整理内容"},
  step_generate:{en:"Generating Answer",zh:"生成回答"},
  step_done:{en:"Completed",zh:"已完成"},
  step_sources:{en:"{0} docs",zh:"{0} 篇文档"},
  ask_empty:{en:"Ask MINISFORUM KB...",zh:"有问题尽管问..."},
  history_title:{en:"History",zh:"历史问答"},
  loading:{en:"Loading...",zh:"加载中..."},
  history_empty:{en:"No history yet",zh:"暂无历史记录"},
  search_title:{en:"Full-text Search",zh:"全文搜索"},
  search_placeholder:{en:"Enter keywords, space-separated",zh:"输入关键词，空格分隔多词"},
  search_btn:{en:"Search",zh:"搜索"},
  search_searching:{en:"Searching...",zh:"搜索中..."},
  search_empty:{en:"No results",zh:"无结果"},
  search_results:{en:"Found {0} results",zh:"找到 {0} 条结果"},
  docs_title:{en:"Document List",zh:"文档列表"},
  docs_upload:{en:"+ Upload",zh:"+ 上传"},
  docs_export:{en:"Export",zh:"导出"},
  docs_empty:{en:"No documents",zh:"暂无文档"},
  docs_diff_title:{en:"Document Diff",zh:"文档差异对比"},
  docs_diff_desc:{en:"Upload old and new versions to highlight changes",zh:"上传旧版和新版两个文件，自动标出变更的文本行"},
  docs_diff_old:{en:"Old file",zh:"旧文件"},
  docs_diff_new:{en:"New file",zh:"新文件"},
  docs_diff_compare:{en:"Compare",zh:"对比"},
  docs_diff_comparing:{en:"Comparing...",zh:"比对中..."},
  docs_diff_pick_both:{en:"Please select both old and new files",zh:"请同时选择旧版和新版两个文件"},
  docs_count:{en:"{0} docs",zh:"{0} 个文档"},
  docs_delete:{en:"Delete",zh:"删除"},
  docs_confirm_delete:{en:"Confirm delete this document?",zh:"确认删除此文档？"},
  docs_deleted:{en:"Deleted",zh:"已删除"},
  docs_delete_failed:{en:"Delete failed: ",zh:"删除失败: "},
  docs_imported:{en:"Imported {0} docs",zh:"导入 {0} 个文档"},
  docs_ocr_hint:{en:"Scanned PDFs & images are auto-OCR'd",zh:"扫描件/图片将自动 OCR 识别入库"},
  docs_upload_failed:{en:"Upload failed: ",zh:"上传失败: "},
  docs_diff_added:{en:"added",zh:"新增"},
  docs_diff_removed:{en:"removed",zh:"删除"},
  docs_diff_line_btn:{en:"Line Diff",zh:"行级对比"},
  spec_compare_title:{en:"Spec Compare",zh:"规格对比"},
  spec_compare_desc:{en:"Upload 2+ product spec files for structured comparison",zh:"上传2个以上产品规格文件，自动提取参数生成对比表格"},
  spec_compare_files:{en:"Select files",zh:"选择文件"},
  spec_compare_btn:{en:"Spec Compare",zh:"规格对比"},
  spec_compare_comparing:{en:"Analyzing...",zh:"分析中..."},
  spec_compare_pick_multi:{en:"Select at least 2 files (max 10)",zh:"至少选择2个文件（最多10个）"},
  spec_compare_mode_spec:{en:"SpecMD Parsed",zh:"规格解析"},
  spec_compare_mode_llm:{en:"AI Extracted",zh:"AI 提取"},
  spec_compare_analysis:{en:"Analysis",zh:"总结分析"},
  spec_compare_no_field:{en:"—",zh:"—"},
  docs_exported:{en:"Exported {0} docs",zh:"已导出 {0} 个文档"},
  docs_export_failed:{en:"Export failed",zh:"导出失败"},
  docs_uncategorized:{en:"Uncategorized",zh:"未分类"},
  error_prefix:{en:"Error: ",zh:"错误: "},
  queue_waiting:{en:"{0} request(s) ahead, please wait...",zh:"前面还有 {0} 个请求在排队，请稍候..."},
  queue_your_turn:{en:"It's your turn, processing...",zh:"轮到您了，正在处理..."},
  generating_in_progress:{en:"Answer is being generated, please wait...",zh:"正在生成回答，请稍候..."},
    empty_answer:{en:"Answer generation failed, please retry.",zh:"回答生成失败，请重试。"},
  empty_answer_hint:{en:"If the problem persists, try rephrasing your question with more details.",zh:"如果问题持续出现，可以尝试补充更多细节后重新提问。"},
  generating:{en:"Generating...",zh:"生成中..."}
};
let _lang = localStorage.getItem('kb_lang') || 'en';
function t(key, ...args) {
  let entry = _i18n[key];
  let s = entry ? (entry[_lang] || entry['en']) : key;
  for(let i=0; i<args.length; i++) s = s.replace('{'+i+'}', args[i]);
  return s;
}
function setLang(l) {
  _lang = l;
  localStorage.setItem('kb_lang', l);
  let sw = document.getElementById('lang-switch'); if(sw) sw.value = l;
  applyLang();
  let activeTab = document.querySelector('.nav-tab.active');
  if(activeTab) switchPage(activeTab.dataset.page || 'ask');
}
function applyLang() {
  let el;
  el = document.getElementById('login-username'); if(el) el.placeholder = t('username');
  el = document.getElementById('login-password'); if(el) el.placeholder = t('password');
  let loginBtns = document.querySelectorAll('#login-modal button'); loginBtns.forEach(b => { if(b.textContent.trim()==='Sign In'||b.textContent.trim()==='登录') b.textContent = t('sign_in'); });
  let loginHint = document.querySelector('#login-modal p'); if(loginHint) loginHint.textContent = t('login_hint');
  let logoutEl = document.querySelector('.sidebar-user .logout'); if(logoutEl) logoutEl.textContent = t('logout');
  let tabs = document.querySelectorAll('.nav-tab');
  let tabKeys = ['ask','history','docs','search'];
  tabs.forEach((tab,i) => { if(i<tabKeys.length) tab.textContent = t(tabKeys[i]); });
  let stLbls = document.querySelectorAll('.stat-lbl');
  let stKeys = ['stats_docs','stats_chunks','stats_queries','stats_users'];
  stLbls.forEach((lbl,i) => { if(i<stKeys.length) lbl.textContent = t(stKeys[i]); });
  document.querySelectorAll('[data-i18n]').forEach(el => {
    let key = el.dataset.i18n;
    if (_i18n[key]) el.textContent = _i18n[key][_lang];
  });
  el = document.querySelector('.related-q-title'); if(el) el.textContent = t('related_title');
  el = document.getElementById('ask-input'); if(el) el.placeholder = t('input_placeholder');
  let sendBtn = document.querySelector('.input-area button'); if(sendBtn) sendBtn.textContent = t('send');
  // Update data-default on file inputs for language switch consistency
  ['diff-old','diff-new','spec-files'].forEach(id => {
    let inp = document.getElementById(id);
    if (inp) {
      let key = id === 'spec-files' ? 'spec_compare_files' : (id === 'diff-old' ? 'docs_diff_old' : 'docs_diff_new');
      inp.setAttribute('data-default', t(key));
    }
  });
}
// ── End i18n ──

let token='', user=null, askContentCache='', askSourcesCache='', isComposing=false, isGenerating=false;
const API = window.location.origin;

async function api(method, path, body) {
  let opts = {method, headers:{'Content-Type':'application/json'}};
  if(token) opts.headers['Authorization'] = 'Bearer '+token;
  if(body) opts.body = JSON.stringify(body);
  let r = await fetch(API+path, opts);
  if(r.status===401) { doLogout(); throw new Error('Unauthorized'); }
  return r.json();
}

async function doLogin() {
  let u = document.getElementById('login-username').value.trim();
  let p = document.getElementById('login-password').value;
  if(!u||!p) return;
  try {
    let r = await api('POST','/api/login',{username:u,password:p});
    if(r.ok) {
      token=r.token; user=r.user;
      localStorage.setItem('kb_token',token); localStorage.setItem('kb_user',JSON.stringify(user));
      updateUI();
    } else toast(r.error||t('login_failed'),true);
  } catch(e) { toast(t('login_failed'),true); }
}

function doLogout() {
  token=''; user=null;
  localStorage.removeItem('kb_token'); localStorage.removeItem('kb_user');
  updateUI();
}

function updateUI() {
  if(token&&user) {
    document.getElementById('login-modal').style.display='none';
    document.getElementById('sidebar').style.display='flex';
    document.getElementById('main-panel').style.display='flex';
    document.getElementById('user-info').textContent = user.username + ' (' + (user.department||'') + ')';
    let sw = document.getElementById('lang-switch'); if(sw) sw.value = _lang;
    applyLang();
    loadStats();
    loadRelatedQuestions();
  } else {
    document.getElementById('login-modal').style.display='flex';
    document.getElementById('sidebar').style.display='none';
    document.getElementById('main-panel').style.display='none';
  }
}

async function loadStats() {
  try {
    let r = await api('GET','/api/stats');
    if(r.ok) {
      document.getElementById('st-docs').textContent = r.docs;
      document.getElementById('st-chunks').textContent = r.chunks;
      document.getElementById('st-queries').textContent = r.queries_today;
      document.getElementById('st-users').textContent = r.users;
    }
  } catch(e) {}
}

async function loadRelatedQuestions() {
  try {
    let r = await api('GET','/api/related-questions');
    if(r.ok && r.questions) {
      let html = r.questions.map(q => `<span class="related-q-tag" onclick="document.getElementById('ask-input').value='${q.replace(/'/g,"\\'")}';switchPage('ask')">${q}</span>`).join('');
      document.getElementById('related-tags').innerHTML = html;
    }
  } catch(e) {}
}

let _ask_seq = 0;
function setGenerating(state) {
  isGenerating = state;
  let input = document.getElementById('ask-input');
  let btn = document.querySelector('.input-area button');
  if (state) {
    input.disabled = true;
    btn.disabled = true;
    btn.textContent = t('generating');
    btn.classList.add('loading');
  } else {
    input.disabled = false;
    btn.disabled = false;
    btn.textContent = t('send');
    btn.classList.remove('loading');
  }
}
async function askQuestion() {
  let q = document.getElementById('ask-input').value.trim();
  if(!q||!token) return;
  if(isGenerating) { toast(t('generating_in_progress'), true); return; }
  setGenerating(true);
  document.getElementById('ask-input').value = '';

  let seq = ++_ask_seq;
  let area = document.getElementById('chat-area');
  area.innerHTML += `<div class="message user">${q}</div>`;
  let progressHTML = `
    <div class="kb-progress-card" id="load-ind-${seq}">
      <div class="kb-progress-bar-wrap">
        <div class="kb-progress-fill loading"></div>
      </div>
      <div class="kb-progress-meta">
        <span class="kb-progress-status">${t('step_search')}...</span>
      </div>
    </div>`;
  area.innerHTML += `<div class="message assistant" id="stream-msg-${seq}"></div>` + progressHTML;
  area.scrollTop = area.scrollHeight;

  // 假性推进: 0→30→70，真实输出时一把到 100%
  let fakeTimers = [];
  let fakeStep = 0;
  function advanceFake() {
    fakeStep++;
    let card = document.getElementById('load-ind-'+seq);
    if(!card) return;
    if(fakeStep === 1) {
      card.querySelector('.kb-progress-fill').style.width = '30%';
      card.querySelector('.kb-progress-status').textContent = t('step_analyze') + '...';
    } else if(fakeStep === 2) {
      card.querySelector('.kb-progress-fill').style.width = '70%';
      card.querySelector('.kb-progress-status').textContent = t('step_search') + '...';
    }
  }
  fakeTimers.push(setTimeout(advanceFake, 1500));
  fakeTimers.push(setTimeout(advanceFake, 3500));

  function stepFull() {
    fakeTimers.forEach(clearTimeout); fakeTimers = [];
    let card = document.getElementById('load-ind-'+seq);
    if(!card) return;
    let fill = card.querySelector('.kb-progress-fill');
    fill.classList.remove('loading');
    fill.style.width = '100%';
    card.querySelector('.kb-progress-status').textContent = t('step_done');
  }
  function dismissCard() {
    let card = document.getElementById('load-ind-'+seq);
    if(!card) return;
    card.classList.add('done');
    setTimeout(() => { if(card.parentNode) card.remove(); }, 500);
  }

  let streamEl = document.getElementById(`stream-msg-${seq}`);
  let loadEl = document.getElementById(`load-ind-${seq}`);
  let fullText = '';
  let firstToken = true;
  let reasoningText = '';
  let renderTimer = null; // 渲染节流定时器

  try {
    let resp = await fetch(API+'/api/kb/ask', {
      method:'POST',
      headers:{'Content-Type':'application/json','Authorization':'Bearer '+token},
      body:JSON.stringify({question:q, lang:_lang})
    });

    const reader = resp.body.getReader();
    const decoder = new TextDecoder();
    let buffer = '';

    while(true) {
      let {done, value} = await reader.read();
      if(done) break;
      buffer += decoder.decode(value, {stream:true});
      let lines = buffer.split('\n');
      buffer = lines.pop();

      for(let line of lines) {
        if(!line.startsWith('data: ')) continue;
        try {
          let obj = JSON.parse(line.slice(6));
          if(obj.type==='sources') {
            document.getElementById('sources-bar').innerHTML = obj.sources.map(s=>{
              let safeTitle = s.title.replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');
              return `<span class="source-chip" onclick="viewDocument(${s.doc_id},'ask')" title="${safeTitle}"><span class="source-chip-label">${safeTitle}</span><span class="source-chip-score">${Math.round(s.score*100)}%</span></span>`;
            }).join('');
          } else if(obj.type==='queue') {
            let statusEl = document.querySelector('#load-ind-'+seq+' .kb-progress-status');
            if(statusEl) {
              if(obj.i18n_key === 'queue_waiting') {
                statusEl.textContent = t('queue_waiting', obj.n || 0);
              } else if(obj.i18n_key === 'queue_your_turn') {
                statusEl.textContent = t('queue_your_turn');
              }
            }
          } else if(obj.type==='reasoning') {
            reasoningText += obj.content;
            let lines = reasoningText.split('\n').filter(l => l.trim());
            let display = lines.slice(-2).join('\n');
            let rel = streamEl.querySelector('.kb-reasoning');
            if(!rel) {
              rel = document.createElement('div');
              rel.className = 'kb-reasoning';
              streamEl.prepend(rel);
            }
            rel.textContent = display;
            rel.scrollTop = rel.scrollHeight;
          } else if(obj.type==='token') {
            if(firstToken) {
              stepFull();
              dismissCard();
              firstToken = false;
            }
            fullText += obj.content;
            // 节流渲染：合并高频 token，避免每个 token 全量 marked.parse + innerHTML
            if(!renderTimer) {
              renderTimer = setTimeout(() => {
                renderTimer = null;
                streamEl.innerHTML = formatMarkdown(fullText);
                area.scrollTop = area.scrollHeight;
              }, 60);
            }
          } else if(obj.type==='done') {
            setGenerating(false);
            stepFull();
            dismissCard();
            if(renderTimer) { clearTimeout(renderTimer); renderTimer = null; }
            streamEl.innerHTML = formatMarkdown(fullText);
            area.scrollTop = area.scrollHeight;
            if(!fullText) {
              streamEl.innerHTML = `<p style="color:var(--warning)">${t('empty_answer') || '回答生成失败，请重试。'}</p><p style="color:var(--muted);font-size:0.9em;margin-top:6px">${t('empty_answer_hint') || '如果问题持续出现，可以尝试补充更多细节后重新提问。'}</p>`;
            }
            askContentCache = area.innerHTML;
            let sb = document.getElementById('sources-bar');
            if(sb) askSourcesCache = sb.innerHTML;
          } else if(obj.type==='error') {
            setGenerating(false);
            if(renderTimer) { clearTimeout(renderTimer); renderTimer = null; }
            streamEl.innerHTML += `<p style="color:var(--danger)">${t('error_prefix')}${obj.content}</p>`;
            if(loadEl) loadEl.remove();
          }
        } catch(e) {}
      }
    }
    // 兜底：流结束但未收到 done/error 事件
    if(isGenerating) {
      setGenerating(false);
      stepFull();
      dismissCard();
      if(renderTimer) { clearTimeout(renderTimer); renderTimer = null; }
      streamEl.innerHTML = formatMarkdown(fullText);
      area.scrollTop = area.scrollHeight;
      if(!fullText) {
        streamEl.innerHTML = `<p style="color:var(--warning)">${t('empty_answer') || '回答生成失败，请重试。'}</p><p style="color:var(--muted);font-size:0.9em;margin-top:6px">${t('empty_answer_hint') || '如果问题持续出现，可以尝试补充更多细节后重新提问。'}</p>`;
      }
    }
  } catch(e) {
    setGenerating(false);
    if(renderTimer) { clearTimeout(renderTimer); renderTimer = null; }
    if(loadEl) loadEl.remove();
    streamEl.innerHTML = `<p style="color:var(--danger)">${t('error_prefix')}${e.message}</p>`;
  }
}

function formatMarkdown(text) {
  if (typeof marked === 'undefined') {
    // CDN 加载失败时的兜底
    return text.replace(/\n/g, '<br>');
  }
  marked.setOptions({ breaks: true, gfm: true });
  return marked.parse(text);
}

function switchPage(pg) {
  let area = document.getElementById('chat-area');
  // Save ask content before switching away
  let curTab = document.querySelector('.nav-tab.active');
  if(curTab && curTab.dataset.page === 'ask' && area.innerHTML && !area.querySelector('p[style*="text-align:center"]') && !area.querySelector('button[onclick*="switchPage"]')) {
    askContentCache = area.innerHTML;
    let sb = document.getElementById('sources-bar');
    if(sb) askSourcesCache = sb.innerHTML;
  }
  document.querySelectorAll('.nav-tab').forEach(el=>el.classList.toggle('active', el.dataset.page===pg));
  document.querySelectorAll('.feat-item').forEach(el=>el.classList.toggle('active', el.dataset.page===pg));
  let bar = document.getElementById('sources-bar');
  if(bar) bar.innerHTML = '';

  if(pg==='ask') {
    area.innerHTML = askContentCache || '<p style="color:var(--text-dim);padding:40px;text-align:center">'+t('ask_empty')+'</p>';
    if(bar && askSourcesCache) bar.innerHTML = askSourcesCache;
    return;
  }
  if(pg==='history') { renderHistory(area); return; }
  if(pg==='docs') { renderDocs(area); return; }
  if(pg==='search') { renderSearch(area); return; }
  if(pg==='compare') { renderCompare(area); return; }
}

async function renderHistory(area) {
  area.innerHTML = `<div style="padding:24px 28px">
    <h3 style="font-size:16px;font-weight:600;margin-bottom:16px">${t('history_title')}</h3>
    <div id="history-list"><p style="color:var(--text-dim)">${t('loading')}</p></div>
  </div>`;
  try {
    let r = await api('GET','/api/history');
    if(!r.ok || !r.history || !r.history.length) {
      document.getElementById('history-list').innerHTML='<p style="color:var(--text-dim)">'+t('history_empty')+'</p>';
      return;
    }
    let html = '';
    for(let h of r.history) {
      html += `<div style="margin-bottom:14px;padding:14px;background:var(--surface);border-radius:var(--radius);border:1px solid var(--glass-border);cursor:pointer"
        onclick="document.getElementById('ask-input').value='${h.question.replace(/['"]/g,'').slice(0,200)}';switchPage('ask')">
        <div style="font-weight:500;font-size:13px;margin-bottom:6px;color:var(--text);word-break:break-word">${h.question.slice(0,200)}</div>
        <div style="font-size:12px;color:var(--text-muted);display:-webkit-box;-webkit-line-clamp:2;-webkit-box-orient:vertical;overflow:hidden">${h.answer.slice(0,300)}</div>
        <div style="font-size:10px;color:var(--text-dim);margin-top:6px">${(h.created_at||'').slice(0,16)}</div>
      </div>`;
    }
    document.getElementById('history-list').innerHTML = html;
  } catch(e) { document.getElementById('history-list').innerHTML='<p style="color:var(--danger)">'+t('error_prefix')+e.message+'</p>'; }
}

async function renderSearch(area) {
  area.innerHTML = `<div style="padding:24px 28px">
    <h3 style="font-size:16px;font-weight:600;margin-bottom:16px">${t('search_title')}</h3>
    <div style="display:flex;gap:8px;margin-bottom:20px">
      <input id="search-input" placeholder="${t('search_placeholder')}" style="flex:1;padding:10px 14px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius);color:var(--text);font:inherit"
        onkeydown="if(event.key==='Enter')doSearch()">
      <button onclick="doSearch()" style="padding:10px 20px;background:var(--primary);border:none;border-radius:var(--radius);color:#fff;font:inherit;font-weight:600;cursor:pointer">${t('search_btn')}</button>
    </div>
    <div id="search-results"></div>
  </div>`;
}

async function doSearch() {
  let q = document.getElementById('search-input').value.trim();
  if(!q) return;
  let el = document.getElementById('search-results');
  el.innerHTML = '<p style="color:var(--text-dim)">'+t('search_searching')+'</p>';
  try {
    let r = await api('GET','/api/search?q='+encodeURIComponent(q)+'&lang='+_lang);
    if(!r.ok || !r.results || !r.results.length) {
      el.innerHTML = '<p style="color:var(--text-dim)">'+t('search_empty')+'</p>';
      return;
    }
    let html = '<p style="font-size:12px;color:var(--text-muted);margin-bottom:12px">'+t('search_results', r.results.length)+'</p>';
    for(let item of r.results) {
      let fromPage = (document.querySelector('.nav-tab.active')||{dataset:{}}).dataset.page || 'search';
      html += `<div style="margin-bottom:14px;padding:14px;background:var(--surface);border-radius:var(--radius);border:1px solid var(--glass-border);cursor:pointer" onclick="viewDocument(${item.doc_id},'${fromPage}')" onmouseover="this.style.borderColor='var(--primary)'" onmouseout="this.style.borderColor='var(--glass-border)'">
        <div style="font-weight:500;font-size:13px;margin-bottom:6px;color:var(--primary)">${item.doc_title} <span style="font-weight:400;color:var(--text-muted);font-size:11px">${item.filename||''}</span></div>
        <div style="font-size:12px;color:var(--text);line-height:1.5">${item.snippet}</div>
      </div>`;
    }
    el.innerHTML = html;
  } catch(e) { el.innerHTML = '<p style="color:var(--danger)">'+t('error_prefix')+e.message+'</p>'; }
}

async function viewDocument(docId, fromPage) {
  fromPage = fromPage || 'search';
  try {
    let r = await api('GET','/api/docs/'+docId);
    if(r.ok && r.doc) {
      let d = r.doc;
      let chatArea = document.getElementById('chat-area');
      chatArea.innerHTML = `<div style="padding:24px 28px">
        <div style="display:flex;align-items:center;gap:12px;margin-bottom:20px">
          <button onclick="switchPage('${fromPage}')" style="padding:6px 12px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius-pill);color:var(--text-dim);font-size:12px;cursor:pointer">&larr; ${t(fromPage)}</button>
          <h3 style="font-size:16px;font-weight:600;margin:0">${d.title}</h3>
          <span style="font-size:11px;color:var(--text-muted);margin-left:auto">${d.filename||''}</span>
        </div>
        <div style="display:flex;gap:12px;margin-bottom:16px">
          ${d.filepath ? `<button onclick="openDocFile(${docId})" style="padding:8px 16px;background:var(--primary);color:#fff;border:none;border-radius:var(--radius);font-size:12px;cursor:pointer;display:inline-flex;align-items:center;gap:6px;font-weight:600">📄 查看原始文件</button>` : ''}
          <span style="font-size:12px;color:var(--text-muted);line-height:36px">${(d.created_at||'').slice(0,10)}</span>
        </div>
        <div style="background:var(--surface);border-radius:var(--radius);border:1px solid var(--glass-border);padding:20px;max-height:70vh;overflow:auto;font-size:13px;line-height:1.8;white-space:pre-wrap">${d.content||'(no text content)'}</div>
      </div>`;
      document.querySelectorAll('.nav-tab').forEach(el=>el.classList.toggle('active', el.dataset.page===fromPage));
      document.querySelectorAll('.feat-item').forEach(el=>el.classList.toggle('active', el.dataset.page===fromPage));
      document.getElementById('sources-bar').innerHTML = '';
    } else {
      toast(t('error_prefix') + 'Document not found', true);
    }
  } catch(e) {
    toast(t('error_prefix') + e.message, true);
  }
}

async function openDocFile(docId) {
  try {
    let resp = await fetch(API+'/api/docs/'+docId+'/file', {
      headers:{'Authorization':'Bearer '+token}
    });
    if(!resp.ok) {
      let err = await resp.json().catch(()=>({detail:'Request failed'}));
      toast(t('error_prefix') + (err.detail||'Download failed'), true);
      return;
    }
    let blob = await resp.blob();
    let url = URL.createObjectURL(blob);
    window.open(url, '_blank');
    setTimeout(()=>URL.revokeObjectURL(url), 60000);
  } catch(e) {
    toast(t('error_prefix') + e.message, true);
  }
}

async function renderDocs(area) {
  area.innerHTML = `<div style="padding:24px 28px">
    <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:16px">
      <h3 style="font-size:16px;font-weight:600">${t('docs_title')}</h3>
      <div style="display:flex;gap:8px">
        <input type="file" id="doc-upload-input" multiple accept=".pdf,.docx,.xlsx,.csv,.pptx,.txt,.md,.html,.htm,.json,.png,.jpg,.jpeg,.bmp,.webp,.tiff" style="display:none" onchange="doUpload()">
        <button onclick="document.getElementById('doc-upload-input').click()" style="padding:8px 16px;background:var(--primary);border:none;border-radius:var(--radius);color:#fff;font:inherit;font-weight:600;cursor:pointer;font-size:12px">${t('docs_upload')}</button>
        <button onclick="doExport()" style="padding:8px 16px;background:transparent;border:1px solid var(--glass-border);border-radius:var(--radius);color:var(--text-dim);font:inherit;cursor:pointer;font-size:12px">${t('docs_export')}</button>
      </div>
    </div>
    <p style="font-size:11px;color:var(--text-muted);margin:-8px 0 12px">${t('docs_ocr_hint')}</p>
    <div id="docs-list"><p style="color:var(--text-dim)">${t('loading')}</p></div>
    <div style="margin-top:28px;padding-top:20px;border-top:1px solid var(--glass-border)">
      <h4 style="font-size:14px;font-weight:600;margin-bottom:12px">${t('docs_diff_title')}</h4>
      <p style="font-size:11px;color:var(--text-muted);margin-bottom:12px">${t('docs_diff_desc')}</p>
      <div style="display:flex;gap:8px;align-items:center;flex-wrap:wrap;margin-bottom:12px">
        <label for="diff-old" style="padding:6px 12px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius-pill);font-size:11px;cursor:pointer;color:var(--text-dim)">
          <span id="diff-old-label">${t('docs_diff_old')}</span>
        </label>
        <input type="file" id="diff-old" style="display:none" onchange="var n=this.files[0]?.name;document.getElementById('diff-old-label').textContent=n||document.querySelector('#diff-old-label').getAttribute('data-default')" data-default="${t('docs_diff_old')}">
        <label for="diff-new" style="padding:6px 12px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius-pill);font-size:11px;cursor:pointer;color:var(--text-dim)">
          <span id="diff-new-label">${t('docs_diff_new')}</span>
        </label>
        <input type="file" id="diff-new" style="display:none" onchange="var n=this.files[0]?.name;document.getElementById('diff-new-label').textContent=n||document.querySelector('#diff-new-label').getAttribute('data-default')" data-default="${t('docs_diff_new')}">
        <button onclick="doDiff()" style="padding:6px 16px;background:var(--primary);border:none;border-radius:var(--radius-pill);color:#fff;font:inherit;font-weight:600;cursor:pointer;font-size:11px">${t('docs_diff_line_btn')}</button>
      </div>
      <div id="diff-results"></div>
      <div style="margin-top:20px;padding-top:16px;border-top:1px solid var(--glass-border)">
        <h4 style="font-size:14px;font-weight:600;margin-bottom:12px">${t('spec_compare_title')}</h4>
        <p style="font-size:11px;color:var(--text-muted);margin-bottom:12px">${t('spec_compare_desc')}</p>
        <div style="display:flex;gap:8px;align-items:center;flex-wrap:wrap;margin-bottom:12px">
          <label for="spec-files" style="padding:6px 12px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius-pill);font-size:11px;cursor:pointer;color:var(--text-dim)">
            <span id="spec-files-label">${t('spec_compare_files')}</span>
          </label>
          <input type="file" id="spec-files" multiple style="display:none" onchange="var n=this.files.length;document.getElementById('spec-files-label').textContent=n?n+' files':document.querySelector('#spec-files-label').getAttribute('data-default')" data-default="${t('spec_compare_files')}">
          <button onclick="doSpecCompare()" style="padding:6px 16px;background:var(--primary);border:none;border-radius:var(--radius-pill);color:#fff;font:inherit;font-weight:600;cursor:pointer;font-size:11px">${t('spec_compare_btn')}</button>
        </div>
        <div id="spec-results"></div>
      </div>
    </div>
  </div>`;
  loadDocs();
}

async function loadDocs() {
  let el = document.getElementById('docs-list');
  try {
    let r = await api('GET','/api/docs?page_size=100');
    if(!r.ok || !r.docs) { el.innerHTML='<p style="color:var(--text-dim)">'+t('docs_empty')+'</p>'; return; }
    let groups = {};
    for(let d of r.docs) {
      let dept = d.department || t('docs_uncategorized');
      if(!groups[dept]) groups[dept] = [];
      groups[dept].push(d);
    }
    let depts = Object.keys(groups).sort();
    let html = '';
    for(let dept of depts) {
      let docs = groups[dept];
      let gid = 'grp-' + dept.replace(/[^a-zA-Z0-9\u3400-\u9fff\uf900-\ufaff]/g, '_');
      html += `<div style="margin-top:16px;margin-bottom:6px;padding:6px 12px;background:var(--surface);border-radius:var(--radius);font-size:12px;font-weight:600;color:var(--primary);display:flex;justify-content:space-between;align-items:center;cursor:pointer" onclick="var b=document.getElementById('${gid}');var a=this.querySelector('.fold-arrow');b.style.display=b.style.display==='none'?'':'none';a.textContent=b.style.display==='none'?'▶':'▼'">
        <span><span class="fold-arrow" style="display:inline-block;width:14px;font-size:10px">▼</span> ${dept}</span><span style="font-weight:400;color:var(--text-dim);font-size:11px">${t('docs_count', docs.length)}</span>
      </div>`;
      html += `<div id="${gid}" style="">`;
      for(let d of docs) {
        let verTag = d.version > 1 ? `<span style="color:var(--primary);font-size:10px;background:rgba(99,102,241,0.12);padding:1px 6px;border-radius:var(--radius-pill);margin-left:4px">v${d.version}</span>` : '';
        html += `<div style="display:flex;justify-content:space-between;align-items:center;padding:10px 0 10px 8px;border-bottom:1px solid var(--glass-border);font-size:13px;min-height:48px;box-sizing:border-box">
          <div style="overflow:hidden;flex:1;min-width:0;padding-right:8px">
            <div style="font-weight:500;white-space:nowrap;overflow:hidden;text-overflow:ellipsis" title="${d.title}">${d.title}${verTag}</div>
            <div style="color:var(--text-muted);font-size:11px">${d.filename || ''} · ${(d.created_at||'').slice(0,10)}</div>
          </div>
          <div style="display:flex;gap:6px">
            <button onclick="showVersions(${d.id})" style="padding:4px 10px;background:transparent;border:1px solid var(--primary);border-radius:var(--radius-pill);color:var(--primary);font-size:11px;cursor:pointer">${t('version_history')}</button>
            <button onclick="doDelete(${d.id})" style="padding:4px 10px;background:transparent;border:1px solid var(--danger);border-radius:var(--radius-pill);color:var(--danger);font-size:11px;cursor:pointer">${t('docs_delete')}</button>
          </div>
        </div>`;
      }
      html += `</div>`;
    }
    el.innerHTML = html || '<p style="color:var(--text-dim)">'+t('docs_empty')+'</p>';
  } catch(e) { el.innerHTML = '<p style="color:var(--danger)">'+t('error_prefix')+e.message+'</p>'; }
}

async function doUpload() {
  let input = document.getElementById('doc-upload-input');
  if(!input.files.length) return;
  let formData = new FormData();
  for(let f of input.files) formData.append('file', f);
  try {
    let resp = await fetch(API+'/api/kb/upload-multipart', {method:'POST',headers:{'Authorization':'Bearer '+token},body:formData});
    let r = await resp.json();
    toast(t('docs_imported', r.imported));
    loadDocs(); loadStats();
  } catch(e) { toast(t('docs_upload_failed')+e.message, true); }
}

async function doDelete(docId) {
  if(!confirm(t('docs_confirm_delete'))) return;
  try {
    await api('DELETE','/api/docs/'+docId);
    toast(t('docs_deleted'));
    loadDocs(); loadStats();
  } catch(e) { toast(t('docs_delete_failed')+e.message, true); }
}

async function doDiff() {
  let oldFile = document.getElementById('diff-old').files[0];
  let newFile = document.getElementById('diff-new').files[0];
  if(!oldFile || !newFile) { toast(t('docs_diff_pick_both'), true); return; }
  let el = document.getElementById('diff-results');
  el.innerHTML = '<p style="color:var(--text-dim)">'+t('docs_diff_comparing')+'</p>';
  let formData = new FormData();
  formData.append('files', oldFile);
  formData.append('files', newFile);
  try {
    let resp = await fetch(API+'/api/docs/diff', {method:'POST',headers:{'Authorization':'Bearer '+token},body:formData});
    let r = await resp.json();
    if(!r.ok) { el.innerHTML = '<p style="color:var(--danger)">'+r.error+'</p>'; return; }
    el.innerHTML = renderDiff(r.diff);
  } catch(e) { el.innerHTML = '<p style="color:var(--danger)">'+t('error_prefix')+e.message+'</p>'; }
}

/* ── Diff 渲染：按 @@ hunk 分组折叠，长行自动换行，支持仅看变更 ── */
function esc(s){return String(s).replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
function toggleHunk(i){let h=document.querySelectorAll('.diff-hunk')[i];if(h)h.classList.toggle('collapsed');}
function setAllHunks(open){document.querySelectorAll('.diff-hunk').forEach(h=>h.classList.toggle('collapsed',!open));}
function toggleDiffContext(show){document.querySelectorAll('.diff-results-box').forEach(b=>b.classList.toggle('diff-hide-ctx',!show));}
function renderDiff(diffText) {
  let hunks = [];
  let cur = null;
  diffText.split('\n').forEach(l => {
    if(l.startsWith('@@')) {
      if(cur) hunks.push(cur);
      cur = {header:l, rows:[], add:0, del:0};
    } else if(cur) {
      cur.rows.push(l);
      if(l.startsWith('+')) cur.add++;
      else if(l.startsWith('-')) cur.del++;
    }
  });
  if(cur) hunks.push(cur);
  let totalAdd = hunks.reduce((s,h)=>s+h.add,0);
  let totalDel = hunks.reduce((s,h)=>s+h.del,0);
  let html = `<div class="diff-toolbar">
    <span style="color:#10b981;font-weight:600">+ ${totalAdd}</span> ${t('docs_diff_added')}
    <span style="color:#ef4444;font-weight:600;margin-left:4px">- ${totalDel}</span> ${t('docs_diff_removed')}
    <span style="flex:1"></span>
    <label style="display:flex;align-items:center;gap:4px;cursor:pointer;color:var(--text-dim);font-size:11px">
      <input type="checkbox" checked onchange="toggleDiffContext(this.checked)" style="accent-color:var(--primary)"> 仅看变更
    </label>
    <button class="btn" onclick="setAllHunks(false)">全部折叠</button>
    <button class="btn" onclick="setAllHunks(true)">全部展开</button>
  </div>
  <div class="diff-results-box diff-hide-ctx">`;
  hunks.forEach((h,i)=>{
    let rowsHtml = '';
    let pendingDel = [];
    h.rows.forEach(l=>{
      if(l.startsWith('+')) {
        if(pendingDel.length) { rowsHtml += diffPairHtml(pendingDel.shift(), l); }
        else { rowsHtml += diffRowHtml(l); }
      } else if(l.startsWith('-')) {
        pendingDel.push(l);
      } else {
        while(pendingDel.length) rowsHtml += diffRowHtml(pendingDel.shift());
        rowsHtml += diffRowHtml(l);
      }
    });
    while(pendingDel.length) rowsHtml += diffRowHtml(pendingDel.shift());
    html += `<div class="diff-hunk" data-idx="${i}">
      <div class="diff-hunk-head" onclick="toggleHunk(${i})" title="点击折叠/展开">
        <span class="diff-hunk-arrow">▾</span>
        <span class="diff-hunk-pos">${esc(h.header)}</span>
        <span class="diff-hunk-badge" style="color:#10b981">+${h.add}</span>
        <span class="diff-hunk-badge" style="color:#ef4444">-${h.del}</span>
      </div>
      <div class="diff-hunk-body">${rowsHtml}</div>
    </div>`;
  });
  html += `</div>`;
  return html;
}
/* 去掉行首 markdown 标题符号（#），标题文字加粗显示 */
function cleanLine(l) {
  let body = l.replace(/^[+-]/, '');
  let m = body.match(/^(#+)\s*(.*)$/);
  if(m) return {heading:true, text:m[2]};
  return {heading:false, text:body};
}
function diffRowHtml(l) {
  let cls = 'diff-ctx';
  if(l.startsWith('+')) cls = 'diff-add';
  else if(l.startsWith('-')) cls = 'diff-del';
  let c = cleanLine(l);
  let content = c.heading ? `<b>${esc(c.text)||'&nbsp;'}</b>` : (esc(c.text)||'&nbsp;');
  return `<div class="diff-row ${cls}">${content}</div>`;
}
function diffPairHtml(oldL, newL) {
  let oc = cleanLine(oldL), nc = cleanLine(newL);
  let o = oc.heading ? `<b>${esc(oc.text)}</b>` : esc(oc.text);
  let n = nc.heading ? `<b>${esc(nc.text)}</b>` : esc(nc.text);
  return `<div class="diff-row diff-pair"><span class="diff-pair-old">${o}</span><span class="diff-pair-arrow">→</span><span class="diff-pair-new">${n}</span></div>`;
}

async function doSpecCompare() {
  let input = document.getElementById('spec-files');
  if(!input.files || input.files.length < 2) { toast(t('spec_compare_pick_multi'), true); return; }
  if(input.files.length > 10) { toast(t('spec_compare_pick_multi'), true); return; }
  let el = document.getElementById('spec-results');
  el.innerHTML = '<p style="color:var(--text-dim)">'+t('spec_compare_comparing')+'</p>';
  let formData = new FormData();
  for(let f of input.files) formData.append('files', f);
  try {
    let resp = await fetch(API+'/api/spec/compare', {method:'POST',headers:{'Authorization':'Bearer '+token},body:formData});
    let r = await resp.json();
    if(!r.ok) { el.innerHTML = '<p style="color:var(--danger)">'+r.error+'</p>'; return; }
    let modeLabel = r.mode === 'spec' ? t('spec_compare_mode_spec') : t('spec_compare_mode_llm');
    let modeColor = r.mode === 'spec' ? '#10b981' : '#f59e0b';
    let html = '<div style="margin-bottom:12px;padding:8px 14px;background:var(--surface);border-radius:var(--radius);font-size:12px">'
      + '<span style="color:'+modeColor+';font-weight:600">'+modeLabel+'</span>'
      + '</div>';
    // Render single flat comparison table
    let table = r.table;
    if(table.rows && table.rows.length) {
      // Use table-layout:fixed + word-break, narrow key column + wide value columns
      let pCount = table.products.length;
      let keyW = pCount <= 2 ? 22 : (pCount <= 4 ? 18 : 14);
      let valW = Math.floor((100 - keyW) / pCount);
      html += '<table style="width:100%;border-collapse:collapse;table-layout:fixed;font-size:11px;border:1px solid var(--glass-border);border-radius:var(--radius);overflow:hidden">';
      html += '<colgroup>';
      html += '<col style="width:'+keyW+'%">';
      for(let ci=0; ci<pCount-1; ci++) html += '<col style="width:'+valW+'%">';
      html += '<col>';
      html += '</colgroup>';
      html += '<thead><tr style="background:var(--surface)"><th style="padding:6px 10px;text-align:left;border-bottom:1px solid var(--glass-border);color:var(--text-dim);word-wrap:break-word;overflow-wrap:break-word;word-break:break-all">'+t('spec_compare_no_field')+'</th>';
      for(let pn of table.products) {
        html += '<th style="padding:6px 10px;text-align:left;border-bottom:1px solid var(--glass-border);color:var(--text-dim);word-wrap:break-word;overflow-wrap:break-word;word-break:break-all">'+pn+'</th>';
      }
      html += '</tr></thead><tbody>';
      for(let row of table.rows) {
        html += '<tr>';
        html += '<td style="padding:6px 10px;border-bottom:1px solid var(--glass-border);font-weight:500;word-wrap:break-word;overflow-wrap:break-word;word-break:break-all">'+row.key+'</td>';
        for(let v of row.values) {
          html += '<td style="padding:6px 10px;border-bottom:1px solid var(--glass-border);word-wrap:break-word;overflow-wrap:break-word;word-break:break-all">'+v+'</td>';
        }
        html += '</tr>';
      }
      html += '</tbody></table>';
    }
    if(r.analysis) {
      html += '<div style="margin-top:16px;padding:14px;background:var(--surface);border-radius:var(--radius);border-left:3px solid var(--primary)">';
      html += '<div style="font-size:12px;font-weight:600;margin-bottom:6px;color:var(--text-dim)">'+t('spec_compare_analysis')+'</div>';
      html += '<div style="font-size:12px;color:var(--text);line-height:1.7">'+r.analysis+'</div>';
      html += '</div>';
    }
    el.innerHTML = html;
  } catch(e) { el.innerHTML = '<p style="color:var(--danger)">'+t('error_prefix')+e.message+'</p>'; }
}

async function doExport() {
  try {
    let r = await api('GET','/api/docs/export');
    if(!r.ok) { toast(t('docs_export_failed'), true); return; }
    let blob = new Blob([JSON.stringify(r, null, 2)], {type:'application/json'});
    let url = URL.createObjectURL(blob);
    let a = document.createElement('a');
    a.href = url;
    a.download = 'kb_export_' + new Date().toISOString().slice(0,10) + '.json';
    a.click();
    URL.revokeObjectURL(url);
    toast(t('docs_exported', r.total_docs));
  } catch(e) { toast(t('docs_export_failed')+': '+e.message, true); }
}

// ── Document Comparison ──
async function renderCompare(area) {
  area.innerHTML = `<div style="padding:24px 28px">
    <h3 style="font-size:16px;font-weight:600;margin-bottom:8px">${t('compare_title')}</h3>
    <p style="font-size:11px;color:var(--text-muted);margin-bottom:16px">${t('compare_desc')}</p>
    <div id="compare-doc-select" style="margin-bottom:12px;max-height:200px;overflow-y:auto;padding:8px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius)">
      <p style="color:var(--text-dim);font-size:12px">${t('loading')}</p>
    </div>
    <div style="display:flex;gap:8px;align-items:center;margin-bottom:16px">
      <button onclick="document.querySelectorAll('.compare-chk').forEach(c=>c.checked=true);document.querySelectorAll('.compare-chk').forEach(c=>c.onchange?.())" style="padding:4px 12px;background:transparent;border:1px solid var(--glass-border);border-radius:var(--radius-pill);color:var(--text-dim);font-size:11px;cursor:pointer">Select All</button>
      <button onclick="document.querySelectorAll('.compare-chk').forEach(c=>c.checked=false);document.querySelectorAll('.compare-chk').forEach(c=>c.onchange?.())" style="padding:4px 12px;background:transparent;border:1px solid var(--glass-border);border-radius:var(--radius-pill);color:var(--text-dim);font-size:11px;cursor:pointer">Deselect All</button>
      <span id="compare-sel-count" style="font-size:11px;color:var(--text-dim);margin-left:4px">0 selected</span>
    </div>
    <div style="display:flex;gap:8px;margin-bottom:16px">
      <input id="compare-dimension" placeholder="${t('compare_dimension')}" style="flex:1;padding:10px 14px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius);color:var(--text);font:inherit;outline:none" onkeydown="if(event.key==='Enter')doCompare()">
      <button onclick="doCompare()" style="padding:10px 20px;background:var(--primary);border:none;border-radius:var(--radius);color:#fff;font:inherit;font-weight:600;cursor:pointer;white-space:nowrap">${t('compare_btn')}</button>
    </div>
    <div id="compare-results"></div>
  </div>`;
  // Load doc list for checkboxes
  try {
    let r = await api('GET','/api/docs?page_size=100&all=1');
    if(!r.ok || !r.docs) return;
    // Group by department for filtering
    let deptDocs = {};
    for(let d of r.docs) {
      let dept = d.department || t('docs_uncategorized');
      if(!deptDocs[dept]) deptDocs[dept] = [];
      deptDocs[dept].push(d);
    }
    let depts = Object.keys(deptDocs).sort();
    let selDept = depts[0] || '';
    function renderCompareChecks(dept) {
      let docs = deptDocs[dept] || [];
      let h = '';
      for(let d of docs) {
        h += `<label style="display:flex;align-items:center;gap:6px;padding:4px 0;font-size:12px;cursor:pointer;color:var(--text)">
          <input type="checkbox" value="${d.id}" class="compare-chk" style="accent-color:var(--primary)">
          <span>${d.title}</span> <span style="color:var(--text-dim);font-size:10px">v${d.version||1}</span>
        </label>`;
      }
      document.getElementById('compare-doc-select').innerHTML = h || '<p style="color:var(--text-dim)">'+t('docs_empty')+'</p>';
      document.querySelectorAll('.compare-chk').forEach(c => c.onchange = () => {
        let n = document.querySelectorAll('.compare-chk:checked').length;
        document.getElementById('compare-sel-count').textContent = n + ' selected';
      });
    }
    // Custom department dropdown
    let curDept = depts[0] || '';
    let deptDropdownHTML = `
      <div id="compare-dept-dropdown" style="position:relative;margin-bottom:8px">
        <div id="compare-dept-toggle" style="padding:8px 30px 8px 12px;background:var(--surface);border:1px solid var(--glass-border);border-radius:var(--radius);color:var(--text);font-size:12px;cursor:pointer;position:relative;user-select:none">
          <span id="compare-dept-label">${curDept}</span>
          <span id="compare-dept-arrow" style="position:absolute;right:10px;top:50%;transform:translateY(-50%);font-size:10px;color:var(--text-dim);transition:transform 0.2s">▼</span>
        </div>
        <div id="compare-dept-pane" style="display:none;position:absolute;top:100%;left:0;right:0;z-index:100;background:var(--bg);border:1px solid var(--glass-border);border-radius:0 0 var(--radius) var(--radius);max-height:180px;overflow-y:auto;box-shadow:0 8px 24px rgba(0,0,0,0.4)">
          ${depts.map(d => `<div class="dept-opt" style="padding:8px 12px;font-size:12px;color:var(--text-dim);cursor:pointer;transition:all 0.15s"
               data-dept="${d.replace(/"/g,'&quot;')}">${d}</div>`).join('')}
        </div>
      </div>
      <div id="compare-doc-select" style="max-height:180px;overflow-y:auto;padding:0 4px"></div>`;
    document.getElementById('compare-doc-select').outerHTML = deptDropdownHTML;
    // Toggle dropdown
    document.getElementById('compare-dept-toggle').addEventListener('click', function() {
      let pane = document.getElementById('compare-dept-pane');
      let arrow = document.getElementById('compare-dept-arrow');
      if(pane.style.display === 'block') {
        pane.style.display = 'none'; arrow.style.transform = '';
      } else {
        pane.style.display = 'block'; arrow.style.transform = 'rotate(180deg)';
      }
    });
    // Option clicks via delegation
    document.getElementById('compare-dept-pane').addEventListener('click', function(e) {
      let opt = e.target.closest('.dept-opt');
      if(!opt) return;
      let dept = opt.dataset.dept;
      renderCompareChecks(dept);
      document.getElementById('compare-dept-label').textContent = dept;
      document.getElementById('compare-dept-pane').style.display = 'none';
      document.getElementById('compare-dept-arrow').style.transform = '';
    });
    // Hover effects via delegation
    document.getElementById('compare-dept-pane').addEventListener('mouseover', function(e) {
      let opt = e.target.closest('.dept-opt');
      if(opt) { opt.style.background = 'rgba(99,102,241,0.12)'; opt.style.color = 'var(--text)'; }
    });
    document.getElementById('compare-dept-pane').addEventListener('mouseout', function(e) {
      let opt = e.target.closest('.dept-opt');
      if(opt) { opt.style.background = ''; opt.style.color = 'var(--text-dim)'; }
    });
    // Close dropdown on outside click
    document.addEventListener('click', function(e) {
      let dd = document.getElementById('compare-dept-dropdown');
      if(dd && !dd.contains(e.target)) {
        let pane = document.getElementById('compare-dept-pane');
        let arrow = document.getElementById('compare-dept-arrow');
        if(pane) { pane.style.display = 'none'; arrow.style.transform = ''; }
      }
    });
    renderCompareChecks(curDept);
  } catch(e) {}
}

function renderCompareText(text) {
  let blocks = text.split(/\n{2,}/);
  let out = '';
  for(let block of blocks) {
    block = block.trim();
    if(!block) continue;
    let lines = block.split('\n');
    // Detect table: 2+ lines, all starting with |
    if(lines.length >= 2 && lines.every(l => l.trim().startsWith('|'))) {
      let headerHtml = '', bodyHtml = '';
      for(let line of lines) {
        let cells = line.split('|').filter(c => c.trim());
        if(!cells.length) continue;
        if(/^[\s\-:]+$/.test(cells.join(''))) continue; // separator row
        let row = '<tr>' + cells.map(c => `<td style="padding:6px 10px;border-bottom:1px solid var(--glass-border);font-size:12px;vertical-align:top">${c.trim().replace(/\*\*(.+?)\*\*/g,'<strong>$1</strong>')}</td>`).join('') + '</tr>';
        if(!headerHtml) headerHtml = row.replace(/<td style="/g,'<th style="background:var(--surface);font-weight:600;').replace(/<\/td>/g,'</th>');
        else bodyHtml += row;
      }
      out += `<div style="overflow-x:auto;margin:16px 0"><table style="width:100%;border-collapse:collapse;border:1px solid var(--glass-border);border-radius:6px">${headerHtml}${bodyHtml}</table></div>`;
      continue;
    }
    // Regular text block
    let html = block
      .replace(/^### (.+)$/gm, '<h4 style="margin-top:16px;font-weight:600;color:var(--text)">$1</h4>')
      .replace(/^## (.+)$/gm, '<h3 style="margin-top:18px;font-weight:700;color:var(--text)">$1</h3>')
      .replace(/^# (.+)$/gm, '<h2 style="margin-top:20px;font-weight:700;color:var(--text)">$1</h2>')
      .replace(/\*\*(.+?)\*\*/g, '<strong>$1</strong>')
      .replace(/^- (.+)$/gm, '<li style="margin-left:16px">$1</li>')
      .replace(/^(\d+)\. (.+)$/gm, '<li style="margin-left:16px">$2</li>')
      .replace(/\n/g, '<br>');
    out += '<div style="margin-bottom:8px;line-height:1.8">' + html + '</div>';
  }
  return '<div style="font-size:13px;color:var(--text);line-height:1.6;padding:8px 0">' + out + '</div>';
}

async function doCompare() {
  let checked = document.querySelectorAll('.compare-chk:checked');
  let docIds = Array.from(checked).map(c => parseInt(c.value));
  let dim = document.getElementById('compare-dimension').value.trim();
  if(docIds.length < 2) { toast(t('compare_no_docs'), true); return; }
  if(!dim) { toast(t('compare_no_dim'), true); return; }
  let el = document.getElementById('compare-results');
  el.innerHTML = '<div style="padding:24px;text-align:center"><div style="display:inline-block;width:24px;height:24px;border:3px solid var(--glass-border);border-top-color:var(--primary);border-radius:50%;animation:spin 0.8s linear infinite;margin-bottom:8px"></div><p style="color:var(--text-dim);font-size:13px">对比生成中，请稍候...</p></div>';
  let acc = '';
  try {
    let resp = await fetch(API+'/api/kb/compare_stream', {
      method: 'POST',
      headers: {'Content-Type':'application/json', 'Authorization':'Bearer '+token},
      body: JSON.stringify({question: dim, dimension: dim, doc_ids: docIds})
    });
    if(!resp.ok || !resp.body) { el.innerHTML = '<p style="color:var(--danger);padding:24px">'+t('compare_no_docs')+'</p>'; return; }
    const reader = resp.body.getReader();
    const decoder = new TextDecoder();
    let buffer = '';
    while(true) {
      const {done, value} = await reader.read();
      if(done) break;
      buffer += decoder.decode(value, {stream:true});
      let idx;
      while((idx = buffer.indexOf('\n\n')) >= 0) {
        let chunk = buffer.slice(0, idx); buffer = buffer.slice(idx+2);
        if(!chunk.startsWith('data: ')) continue;
        let msg;
        try { msg = JSON.parse(chunk.slice(6)); } catch(e) { continue; }
        if(msg.type === 'token') { acc += msg.content; el.innerHTML = renderCompareText(acc); }
        else if(msg.type === 'error') { el.innerHTML = '<p style="color:var(--danger);padding:24px">'+msg.message+'</p>'; return; }
        else if(msg.type === 'done') { el.innerHTML = renderCompareText(acc); el.scrollIntoView({behavior:'smooth',block:'start'}); }
      }
    }
  } catch(e) { el.innerHTML = '<p style="color:var(--danger);padding:24px">'+t('error_prefix')+e.message+'</p>'; }
}
// ── Version Management ──
async function showVersions(docId) {
  let area = document.getElementById('chat-area');
  area.innerHTML = '<p style="color:var(--text-dim);padding:40px;text-align:center">'+t('loading')+'</p>';
  try {
    let r = await api('GET','/api/docs/'+docId+'/versions');
    if(!r.ok || !r.versions) { area.innerHTML = '<p style="color:var(--text-dim);padding:40px;text-align:center">No versions found</p>'; return; }
    let vers = r.versions;
    let html = `<div style="padding:24px 28px">
      <h3 style="font-size:16px;font-weight:600;margin-bottom:4px">${vers[0].title} — ${t('version_history')}</h3>
      <p style="font-size:11px;color:var(--text-muted);margin-bottom:16px">${t('version_diff_pick')}</p>
      <div style="margin-bottom:12px">`;
    for(let v of vers) {
      html += `<label style="display:flex;align-items:center;gap:6px;padding:6px 0;font-size:13px;cursor:pointer;color:var(--text)">
        <input type="checkbox" value="${v.id}" class="ver-diff-chk" style="accent-color:var(--primary)">
        <span style="font-weight:500">v${v.version}</span>
        <span style="color:var(--text-dim);font-size:11px">${(v.created_at||'').slice(0,16)}</span>
        <span style="color:var(--text-dim);font-size:11px">(${v.filename})</span>
      </label>`;
    }
    html += `</div>
      ${vers.length >= 2 ? `<button onclick="doVersionDiff()" style="padding:8px 20px;background:var(--primary);border:none;border-radius:var(--radius);color:#fff;font:inherit;font-weight:600;cursor:pointer;margin-bottom:16px">${t('version_diff_btn')}</button>` : `<p style="font-size:12px;color:var(--text-dim);margin-bottom:16px">${t('version_only_one')}</p>`}
      <button onclick="switchPage('docs')" style="padding:8px 20px;background:transparent;border:1px solid var(--glass-border);border-radius:var(--radius);color:var(--text-dim);font:inherit;cursor:pointer;margin-left:8px">${t('docs_back')}</button>
      <div id="ver-diff-results"></div>
    </div>`;
    area.innerHTML = html;
  } catch(e) { area.innerHTML = '<p style="color:var(--danger)">'+t('error_prefix')+e.message+'</p>'; }
}

async function doVersionDiff() {
  let checked = document.querySelectorAll('.ver-diff-chk:checked');
  let ids = Array.from(checked).map(c => parseInt(c.value));
  if(ids.length !== 2) { toast(t('version_diff_pick'), true); return; }
  let el = document.getElementById('ver-diff-results');
  el.innerHTML = '<p style="color:var(--text-dim)">'+t('docs_diff_comparing')+'</p>';
  try {
    let r = await api('POST','/api/docs/versions/diff', {v1: ids[0], v2: ids[1]});
    if(!r.ok) { el.innerHTML = '<p style="color:var(--danger)">'+r.error+'</p>'; return; }
    el.innerHTML = `<div style="margin-bottom:8px;font-size:11px;color:var(--text-muted)">
      <strong>${r.v1.title} v${r.v1.version}</strong> → <strong>${r.v2.title} v${r.v2.version}</strong>
      </div>` + renderDiff(r.diff);
  } catch(e) { el.innerHTML = '<p style="color:var(--danger)">'+t('error_prefix')+e.message+'</p>'; }
}

function toast(msg, isError) {
  let el = document.createElement('div');
  el.className = 'toast';
  el.textContent = msg;
  if(isError) el.style.borderColor = 'var(--danger)';
  document.body.appendChild(el);
  setTimeout(()=>el.remove(), 3000);
}

(function(){
  // IME composition guard: prevent Enter from sending during Chinese input
  let askInput = document.getElementById('ask-input');
  if(askInput) {
    askInput.addEventListener('compositionstart', ()=>{ isComposing = true; });
    askInput.addEventListener('compositionend', ()=>{ isComposing = false; });
  }
  let t=localStorage.getItem('kb_token'), u=localStorage.getItem('kb_user');
  if(t&&u){ token=t; user=JSON.parse(u);
    api('GET','/api/health').then(r=>{
      if(r.ok) updateUI();
      else { token='';user=null;localStorage.removeItem('kb_token');localStorage.removeItem('kb_user');updateUI(); }
    });
  } else updateUI();
})();
</script>
</body>
</html>"""

# ── Main ─────────────────────────────────────────────────
if __name__ == "__main__":
    print(f"[V10] Effective config: LLM_API_URL={LLM_API_URL!r}  OLLAMA_HOST={OLLAMA_HOST!r}", flush=True)
    print(f"[V10] GPU channel: URL={GPU_LLM_URL!r}  MODEL={GPU_LLM_MODEL!r}  (ask/chat use_gpu=True)", flush=True)
    print(f"[V10] Models — ask:{KB_MODEL_ASK}  chat:{KB_MODEL_CHAT}  compare:{KB_MODEL_COMPARE}  extractor:{KB_MODEL_EXTRACTOR}  embed:{EMBED_MODEL}", flush=True)
    print("[V10] Initializing DB...", flush=True)
    init_db()
    print("[V10] Rebuilding BM25 index...", flush=True)
    bm25.rebuild()
    print("[V10] Scanning NAS documents...", flush=True)
    result = scan_nas()
    print(f"[V10] Scan done — imported: {result['imported']}, skipped: {result['skipped']}, errors: {result['errors']}", flush=True)
    bm25.rebuild()
    print("[V10] Building embedding index...", flush=True)
    _emb_index = None
    _build_embedding_index()
    print("[V10] Pre-warming LLM model...", flush=True)
    try:
        ollama_chat([{"role": "user", "content": "hi"}], model=KB_MODEL_ASK, max_tokens=1)
        print("[V10] LLM model warm", flush=True)
    except Exception as e:
        print(f"[V10] LLM warm-up failed (non-fatal): {e}", flush=True)
    print("[V10] KB Server V10 (FastAPI + SSE) starting on http://0.0.0.0:8080", flush=True)
    uvicorn.run(app, host="0.0.0.0", port=8080, log_level="warning")
